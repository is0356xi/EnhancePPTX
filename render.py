# presentation_generator/render.py
# -*- coding: utf-8 -*-
"""
render.py
- IR(YAML) を読み込み、スライドを決定論的に描画するオーケストレータ。
- 特色:
  * LayoutEngine: anchor 解決（title/left/right, two_panel_bottom 対応）
  * 実グループ化: component に group: true があれば GroupShape 内に描画
  * 論理グルーピング: 追加図形をレジストリに記録（shape_ids / bbox）
  * IR ノーマライザー: list ルートや components 直下など緩い入力も包んで処理
  * ロガー/ツールローディング強化/エラーメッセージ明確化
"""
from __future__ import annotations

import argparse
import importlib
import importlib.util
from pathlib import Path
import re
import sys
import yaml

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from tools import utils  # parse_geom など
from layouts import LayoutEngine  # 新しいLayoutEngine


# =========================================================
# Tool Loader and Schema Validation
# =========================================================
def _sanitize_module_name(name: str) -> str:
    n = (name or "").strip()
    n = n.replace("-", "_").replace(" ", "_")
    n = re.sub(r"[^0-9a-zA-Z_]", "_", n)
    return n.lower()


def load_schema(tool_name: str):
    """ツール名から動的にPydanticスキーマを読み込む"""
    mod_name = _sanitize_module_name(tool_name)
    full = f"tools.schemas.{mod_name}"
    spec = importlib.util.find_spec(full)
    if spec is None:
        print(
            f"Warning: Schema for tool '{tool_name}' not found. "
            f"Expected 'tools/schemas/{mod_name}.py' with 'Schema' class."
        )
        return None
    try:
        module = importlib.import_module(full)
    except Exception as e:
        print(f"Warning: Failed to import schema '{full}': {type(e).__name__}: {e}")
        return None
    if not hasattr(module, "Schema"):
        print(
            f"Warning: Schema module '{full}' loaded, " f"but no 'Schema' class found."
        )
        return None
    return module.Schema


def load_renderer(tool_name: str):
    """ツール名から動的にレンダラーを読み込む"""
    mod_name = _sanitize_module_name(tool_name)
    full = f"tools.renderers.{mod_name}"
    spec = importlib.util.find_spec(full)
    if spec is None:
        print(
            f"Warning: Renderer for tool '{tool_name}' not found. "
            f"Expected 'tools/renderers/{mod_name}.py' with 'render' function."
        )
        return None
    try:
        module = importlib.import_module(full)
    except Exception as e:
        print(f"Warning: Failed to import renderer '{full}': {type(e).__name__}: {e}")
        return None
    if not hasattr(module, "render"):
        print(
            f"Warning: Renderer module '{full}' loaded, "
            f"but no 'render(slide, data, geom, context)' function."
        )
        return None
    return module


def load_tool(tool_name: str):
    """ツール名から動的にモジュールを読み込む（詳細なエラーを出す）- 後方互換性用"""
    mod_name = _sanitize_module_name(tool_name)
    full = f"tools.{mod_name}"
    spec = importlib.util.find_spec(full)
    if spec is None:
        print(
            f"Error: Tool '{tool_name}' not found. "
            f"Ensure 'tools/{mod_name}.py' exists and 'tools/__init__.py' is present."
        )
        return None
    try:
        module = importlib.import_module(full)
    except Exception as e:
        print(f"Error: Failed to import '{full}': {type(e).__name__}: {e}")
        return None
    if not hasattr(module, "render"):
        print(
            f"Error: Tool '{tool_name}' loaded as '{full}', "
            f"but no 'render(slide, data, geom, context)'."
        )
        return None
    return module


# =========================================================
# Slide size
# =========================================================
def set_slide_size(prs: Presentation, size_info: dict):
    """スライドサイズを設定する"""
    if not size_info:
        return
    if "preset" in size_info and size_info["preset"] == "16x9":
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
    elif "w_mm" in size_info and "h_mm" in size_info:
        # mm to EMU (1mm = 36000 EMU)
        prs.slide_width = int(size_info["w_mm"] * 36000)
        prs.slide_height = int(size_info["h_mm"] * 36000)


# =========================================================
# IR normalizer
# =========================================================
def _normalize_ir(ir):
    """
    受け取った IR を {version, meta, theme, slides} に正規化する。
    - list ルート: slides とみなす / components の配列なら 1 スライドに包む
    - dict ルート: slides が無くて components だけなら 1 スライドに包む
    """

    def _mk_min(slides):
        return {"version": 1, "meta": {}, "theme": {}, "slides": slides}

    if isinstance(ir, list):
        if ir and all(isinstance(x, dict) for x in ir):
            # components 配列かスライド配列かをざっくり推測
            if "tool" in ir[0] or "components" in ir[0]:
                if "tool" in ir[0]:
                    # これは components の配列
                    return _mk_min(
                        [{"id": "auto_1", "background": "#FFFFFF", "components": ir}]
                    )
                # スライド配列（とみなす）
                return _mk_min(ir)
        # よく分からない配列 → 1 スライドの components として包む
        return _mk_min([{"id": "auto_1", "background": "#FFFFFF", "components": ir}])

    if isinstance(ir, dict):
        ir.setdefault("version", 1)
        ir.setdefault("meta", {})
        ir.setdefault("theme", {})
        if "slides" not in ir:
            if "components" in ir:
                slide = {
                    "id": ir.get("id", "auto_1"),
                    "background": ir.get("background", "#FFFFFF"),
                    "layout": ir.get("layout"),
                    "components": ir["components"],
                }
                ir["slides"] = [slide]
            else:
                ir["slides"] = []
        return ir

    raise TypeError(f"IR must be dict or list. Got: {type(ir)}")


# =========================================================
# Grouping helpers (logical grouping)
# =========================================================
def _snapshot_ids(slide):
    return {sh.shape_id for sh in slide.shapes}


def _new_shapes(slide, before_ids):
    return [sh for sh in slide.shapes if sh.shape_id not in before_ids]


def _bbox(shapes):
    if not shapes:
        return None
    lefts = [s.left for s in shapes]
    tops = [s.top for s in shapes]
    rights = [s.left + s.width for s in shapes]
    bottoms = [s.top + s.height for s in shapes]
    left = min(lefts)
    top = min(tops)
    width = max(rights) - left
    height = max(bottoms) - top
    return {"left": left, "top": top, "width": width, "height": height}


def _add_group_shape_compat(shapes, left, top, width, height):
    """
    python-pptx のバージョン差を吸収して GroupShape を作成する。
    - 新しめ: add_group_shape(left, top, width, height)
    - 古め  : add_group_shape() のちに .left/.top/.width/.height を設定
    """
    try:
        # 引数あり版（left, top, width, height を渡せる）
        return shapes.add_group_shape(left, top, width, height)
    except TypeError:
        # 引数なし版：作ってから座標・サイズを設定
        grp = shapes.add_group_shape()
        grp.left, grp.top, grp.width, grp.height = left, top, width, height
        return grp


# =========================================================
# Color helpers
# =========================================================
def _force_text_color_on_shapes(shapes_or_list, rgb: RGBColor):
    """
    渡された shapes コレクション or そのリストに対して、
    - テキストボックス／図形の TextFrame（段落・ラン）
    - テーブル内セルの TextFrame
    - グループ図形配下（再帰）
    - チャート（凡例・データラベル・軸ラベル・タイトル）
    の文字色を強制的に rgb で上書き。
    """
    # shapes_or_list は slide.shapes / group.shapes / [shape,...] のいずれも可
    for sh in list(shapes_or_list):
        # グループ（再帰）
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            _force_text_color_on_shapes(sh.shapes, rgb)

        # テキストフレーム
        if getattr(sh, "has_text_frame", False):
            tf = sh.text_frame
            for p in tf.paragraphs:
                # 段落フォント（ランが無い場合にも効く）
                p.font.color.rgb = rgb
                for r in p.runs:
                    r.font.color.rgb = rgb

        # テーブル
        if getattr(sh, "has_table", False):
            for row in sh.table.rows:
                for cell in row.cells:
                    tf = cell.text_frame
                    for p in tf.paragraphs:
                        p.font.color.rgb = rgb
                        for r in p.runs:
                            r.font.color.rgb = rgb

        # チャート
        if sh.shape_type == MSO_SHAPE_TYPE.CHART:
            chart = sh.chart
            try:
                if (
                    chart.has_title
                    and chart.chart_title
                    and chart.chart_title.text_frame
                ):
                    for p in chart.chart_title.text_frame.paragraphs:
                        p.font.color.rgb = rgb
                        for r in p.runs:
                            r.font.color.rgb = rgb
            except Exception:
                pass
            try:
                if chart.has_legend and chart.legend:
                    chart.legend.font.color.rgb = rgb
            except Exception:
                pass
            # 軸ラベル／メモリ
            for ax_attr in ("category_axis", "value_axis", "series_axis"):
                ax = getattr(chart, ax_attr, None)
                if ax is not None and getattr(ax, "tick_labels", None):
                    try:
                        ax.tick_labels.font.color.rgb = rgb
                    except Exception:
                        pass
            # データラベル
            for plot in chart.plots:
                try:
                    if plot.has_data_labels:
                        plot.data_labels.font.color.rgb = rgb
                except Exception:
                    pass


# =========================================================
# Render
# =========================================================
def render_presentation(ir_path: str, output_path: str):
    """YAML IRを読み込み、PowerPointプレゼンテーションを生成する"""
    with open(ir_path, "r", encoding="utf-8") as f:
        raw = yaml.safe_load(f)
    ir = _normalize_ir(raw)

    prs = Presentation()

    # メタデータとサイズ設定
    meta = ir.get("meta", {}) or {}
    if "slide_size" in meta:
        set_slide_size(prs, meta["slide_size"])

    # テーマ設定
    theme = ir.get("theme", {}) or {}
    if "font_color" not in theme:
        theme["font_color"] = "#000000"  # デフォルトの文字色を黒に設定

    # スライドの描画ループ
    for slide_idx, slide_spec in enumerate(ir.get("slides", []), start=1):
        # 空のスライドを追加（テンプレートを使う場合はレイアウトIDを指定）
        blank_slide_layout = prs.slide_layouts[6]  # 6は通常「白紙」
        slide = prs.slides.add_slide(blank_slide_layout)

        # タイトルが存在する場合、自動的にslide_titleコンポーネントを追加
        components = list(slide_spec.get("components", []))
        slide_title = slide_spec.get("title")
        if slide_title:
            # 既存のslide_titleコンポーネントがあるかチェック
            has_slide_title = any(
                comp.get("tool") == "slide_title" for comp in components
            )

            if not has_slide_title:
                # 自動的にslide_titleコンポーネントを先頭に追加
                auto_title_comp = {
                    "tool": "slide_title",
                    "id": "auto_title",
                    "anchor": "title",
                    "z_index": -1000,  # 最背面に配置
                    "data": {"title": slide_title},
                }
                components.insert(0, auto_title_comp)
                print(
                    f"Auto-added slide_title component for slide {slide_idx}: "
                    f"'{slide_title}'"
                )

        # レイアウト解決エンジンの初期化
        resolver = LayoutEngine(slide_spec, prs)

        # コンテキストはスライド内で共有（グループレジストリを持つ）
        class SimpleLogger:
            def info(self, msg):
                print(msg)

            def warning(self, msg):
                print("Warning:", msg)

            def error(self, msg):
                print("Error:", msg)

        context = {
            "prs": prs,
            "theme": theme,
            "logger": SimpleLogger(),
            "registry": {"groups": {}},  # comp_id -> {tool, shape_ids, bbox, is_group}
        }

        # コンポーネントの描画ループ（z_index 昇順）
        # 既に上でタイトル自動追加済みのcomponentsリストを使用
        components = sorted(components, key=lambda c: c.get("z_index", 0))

        for comp_idx, comp in enumerate(components, start=1):
            tool_name = comp.get("tool")
            if not tool_name:
                print(f"Warning: Component is missing 'tool'. Skipping. {comp}")
                continue

            comp_id = comp.get("id") or f"{tool_name}_{slide_idx}_{comp_idx}"

            # FR-3: スキーマの動的ロードと検証
            schema_class = load_schema(tool_name)
            if schema_class is None:
                # スキーマが見つからない場合、後方互換性のため従来のロードを試行
                print(
                    f"Warning: No schema found for '{tool_name}', falling back to legacy tool loading."
                )
                tool = load_tool(tool_name)
                if not tool:
                    print(
                        f"Warning: Tool '{tool_name}' could not be loaded. Skipping component."
                    )
                    continue
                validated_data = comp.get("data", {})
            else:
                # FR-3 & FR-5: スキーマ検証
                try:
                    raw_data = comp.get("data", {})
                    validated_data = schema_class(**raw_data)
                    print(f"Info: Successfully validated data for tool '{tool_name}'")
                except Exception as e:
                    # FR-5: 検証失敗時の処理
                    print(
                        f"Warning: Schema validation failed for component '{comp_id}' (tool: '{tool_name}'): {e}"
                    )
                    print(
                        f"Warning: Skipping component '{comp_id}' due to validation error."
                    )
                    continue

                # FR-2 & FR-4: レンダラーの動的ロード
                renderer_module = load_renderer(tool_name)
                if renderer_module is None:
                    # レンダラーが見つからない場合、後方互換性のため従来のロードを試行
                    print(
                        f"Warning: No renderer found for '{tool_name}', falling back to legacy tool loading."
                    )
                    tool = load_tool(tool_name)
                    if not tool:
                        print(
                            f"Warning: Tool '{tool_name}' could not be loaded. Skipping component."
                        )
                        continue
                else:
                    tool = renderer_module

            try:
                geom = resolver.resolve(comp)
                # テーブルツールの場合、常にgroupをfalseに設定
                if tool_name == "table":
                    use_group = False
                else:
                    use_group = bool(comp.get("group", True))
                if use_group:
                    # 実グループ: GroupShape 内に描画させる
                    grp = _add_group_shape_compat(
                        slide.shapes,
                        geom["left"],
                        geom["top"],
                        geom["width"],
                        geom["height"],
                    )

                    # FR-4: 検証済みデータの受け渡し
                    tool.render(
                        slide=slide,
                        data=validated_data,
                        geom=geom,
                        context={**context, "shapes_target": grp.shapes, "group": grp},
                    )

                    # 文字を強制的に黒に
                    _force_text_color_on_shapes(grp.shapes, RGBColor(0, 0, 0))

                    member_ids = [s.shape_id for s in grp.shapes]
                    context["registry"]["groups"][comp_id] = {
                        "tool": tool_name,
                        "shape_ids": [grp.shape_id] + member_ids,
                        "bbox": {
                            "left": geom["left"],
                            "top": geom["top"],
                            "width": geom["width"],
                            "height": geom["height"],
                        },
                        "is_group": True,
                    }

                else:
                    # 非グループ: slide.shapes に追加された差分を検出して論理グループ化
                    before_ids = _snapshot_ids(slide)
                    tool.render(
                        slide=slide,
                        data=validated_data,
                        geom=geom,
                        context=context,  # ツール側は slide.shapes をそのまま使う
                    )
                    new_shapes = _new_shapes(slide, before_ids)

                    # 文字を強制的に黒に
                    _force_text_color_on_shapes(new_shapes, RGBColor(0, 0, 0))

                    bbox = _bbox(new_shapes)
                    context["registry"]["groups"][comp_id] = {
                        "tool": tool_name,
                        "shape_ids": [s.shape_id for s in new_shapes],
                        "bbox": bbox,
                        "is_group": False,
                    }

            except Exception as e:
                print(f"Error rendering component {comp.get('id', tool_name)}: {e}")
                # 弱い失敗: コンポーネント描画に失敗しても全体処理は継続

    # 保存
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    # print(f"Presentation saved to {output_path}")


# =========================================================
# CLI
# =========================================================
def _build_arg_parser():
    p = argparse.ArgumentParser(description="Render PPTX from IR(YAML).")
    p.add_argument(
        "input", nargs="?", default="examples/sample.yaml", help="Path to IR YAML"
    )
    p.add_argument(
        "-o", "--output", default="dist/output.pptx", help="Path to output .pptx"
    )
    return p


if __name__ == "__main__":
    # パッケージ直下をパスに追加（相対実行の安定化）
    BASE = Path(__file__).resolve().parent
    if str(BASE) not in sys.path:
        sys.path.insert(0, str(BASE))

    args = _build_arg_parser().parse_args()
    render_presentation(args.input, args.output)
