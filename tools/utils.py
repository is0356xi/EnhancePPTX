from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def pct_to_emu(pct, total_emu):
    """パーセンテージをEMUに変換する"""
    return int(total_emu * (pct / 100.0))


def parse_geom(pos_pct, prs):
    """
    パーセンテージで指定された位置とサイズをEMUに変換する。
    pos: {x, y, w, h} in percentage
    prs: Presentation object
    """
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    return {
        "left": pct_to_emu(pos_pct["x"], slide_width),
        "top": pct_to_emu(pos_pct["y"], slide_height),
        "width": pct_to_emu(pos_pct["w"], slide_width),
        "height": pct_to_emu(pos_pct["h"], slide_height),
    }


def hex_to_rgb(hex_color):
    """HEXカラーコードをRGBタプルに変換する"""
    hex_color = hex_color.lstrip("#")
    return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))


def rgb_to_hex(r, g, b):
    """RGBタプルをHEXカラーコードに変換する"""
    return "#{:02X}{:02X}{:02X}".format(
        max(0, min(255, int(r))),
        max(0, min(255, int(g))),
        max(0, min(255, int(b))),
    )


def tint_color(hex_color, factor=0.20):
    """色を明るく（factor>0）/暗く（factor<0）する"""
    hex_color = hex_color.lstrip("#")
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)

    def _blend(c, toward):
        return max(0, min(255, int(round(c + (toward - c) * abs(factor)))))

    if factor >= 0:
        nr, ng, nb = _blend(r, 255), _blend(g, 255), _blend(b, 255)
    else:
        nr, ng, nb = _blend(r, 0), _blend(g, 0), _blend(b, 0)
    return rgb_to_hex(nr, ng, nb)


def mix_with_white(hex_color, ratio=0.70):
    """色を白と混合してパステル調にする"""
    r, g, b = hex_to_rgb(hex_color)
    r = r + (255 - r) * ratio
    g = g + (255 - g) * ratio
    b = b + (255 - b) * ratio
    return rgb_to_hex(r, g, b)


def mix_rgb_with_white(rgb, ratio_white=0.96):
    """RGBタプルを白と線形補間して明度を上げる"""
    r, g, b = rgb
    rw = int(255 * ratio_white + r * (1 - ratio_white))
    gw = int(255 * ratio_white + g * (1 - ratio_white))
    bw = int(255 * ratio_white + b * (1 - ratio_white))
    return (min(255, rw), min(255, gw), min(255, bw))


def get_theme_color(theme, key, fallback):
    """テーマから色を取得してRGBタプルに変換"""
    raw = (theme or {}).get(key, fallback)
    return normalize_color(raw)


def create_text_box(
    shapes,
    left,
    top,
    width,
    height,
    text,
    font_size=10,
    bold=False,
    color="#000000",
    align=PP_ALIGN.LEFT,
    v_anchor=MSO_ANCHOR.MIDDLE,
):
    """共通のテキストボックス作成ヘルパー"""
    tb = shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = v_anchor
    try:
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = 0
    p.space_after = 0
    p.text = str(text).strip()
    f = p.runs[0].font if p.runs else p.font
    f.size = Pt(font_size)
    f.bold = bold
    r, g, b = normalize_color(color)
    f.color.rgb = RGBColor(r, g, b)
    return tb


def create_rounded_shape(shapes, left, top, width, height, bg_color, radius_pct=0.1):
    """角丸の図形を作成（枠線なし）"""
    shape = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    # 角丸設定
    try:
        shape.adjustments[0] = int(max(0.0, min(1.0, radius_pct)) * 50000)
    except Exception:
        pass
    # 塗り
    r, g, b = normalize_color(bg_color)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    shape.line.fill.background()  # 枠線なし
    return shape


def setup_chart_title(chart, title, font_size=16):
    """チャートのタイトルを設定"""
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(font_size)


def setup_chart_legend(chart, show_legend, position, font_size=12):
    """チャートの凡例を設定"""
    if show_legend:
        chart.has_legend = True
        chart.legend.position = position
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(font_size)


def setup_chart_data_labels(chart, data_labels_config, font_size=10):
    """チャートのデータラベルを設定"""
    if data_labels_config and data_labels_config != "none":
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(font_size)
        return data_labels
    return None


def create_chip_with_text(
    shapes,
    left,
    top,
    width,
    height,
    text,
    bg_color,
    radius_pct=0.3,
    font_size=10,
    text_color=None,
):
    """角丸チップ（塗りつぶし図形）にテキストを配置"""
    shape = create_rounded_shape(shapes, left, top, width, height, bg_color, radius_pct)

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # 内側マージン
    try:
        tf.margin_left = tf.margin_right = Pt(8)
        tf.margin_top = tf.margin_bottom = Pt(2)
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_before = 0
    p.space_after = 0
    p.text = str(text).strip()
    f = p.runs[0].font if p.runs else p.font
    f.size = Pt(font_size)
    f.bold = False
    if text_color:
        r, g, b = normalize_color(text_color)
        f.color.rgb = RGBColor(r, g, b)
    else:
        # 背景に対する黒/白の自動判定
        bg_hex = rgb_to_hex(*normalize_color(bg_color))
        f.color.rgb = get_contrast_color(bg_hex)
    return shape


def normalize_color(color_val):
    """HEX(#RRGGBB) or (r,g,b) or [r,g,b] -> (r,g,b)"""
    if isinstance(color_val, (tuple, list)) and len(color_val) == 3:
        return tuple(int(v) for v in color_val)
    if isinstance(color_val, str):
        return hex_to_rgb(color_val)
    return (0, 0, 0)


def get_contrast_color(hex_bg_color):
    """背景色に基づいてコントラストの高い文字色（黒または白）を返す"""
    r, g, b = hex_to_rgb(hex_bg_color)
    # 輝度を計算 (ITU-R BT.709)
    luminance = (0.2126 * r + 0.7152 * g + 0.0722 * b) / 255
    if luminance > 0.5:
        return RGBColor(0, 0, 0)  # 黒
    else:
        return RGBColor(255, 255, 255)  # 白
