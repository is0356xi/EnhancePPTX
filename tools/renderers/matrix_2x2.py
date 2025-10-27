from __future__ import annotations

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from ..utils import normalize_color
from ..schemas.matrix_2x2 import Matrix2x2Schema


def _text_box(
    shapes,
    left,
    top,
    width,
    height,
    text: str,
    pt: int,
    bold=False,
    color="#000000",
    align=PP_ALIGN.CENTER,
    v_anchor=MSO_ANCHOR.MIDDLE,
):
    """テキストボックスを作成"""
    tb = shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = v_anchor
    try:
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = 0
    p.space_after = 0
    p.text = text.strip()
    f = p.runs[0].font if p.runs else p.font
    f.size = Pt(pt)
    f.bold = bold
    r, g, b = normalize_color(color)
    f.color.rgb = RGBColor(r, g, b)
    return tb


def _draw_line(shapes, x1, y1, x2, y2, color, width_pt):
    """直線を描画"""
    line = shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = RGBColor(*normalize_color(color))
    line.line.width = Pt(width_pt)
    return line


def _draw_quadrant_bg(shapes, left, top, width, height, color):
    """象限の背景を描画"""
    rect = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(*normalize_color(color))
    rect.line.fill.background()  # 枠線なし
    return rect


def _draw_circle(shapes, cx, cy, radius, color):
    """円を描画（プロットポイント用）"""
    left = cx - radius
    top = cy - radius
    size = radius * 2
    circle = shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(*normalize_color(color))
    circle.line.fill.background()
    return circle


def render(slide, data: Matrix2x2Schema, geom: dict, context: dict):
    """
    2x2マトリクスを描画。
    geom は EMU 完成形: {"left","top","width","height"}。
    """
    shapes = context.get("shapes_target", slide.shapes)

    # ジオメトリ情報
    container_left = geom["left"]
    container_top = geom["top"]
    container_width = geom["width"]
    container_height = geom["height"]

    # スタイル設定
    style = data.style
    padding_px = int(
        container_width * (style.padding_pct / 100.0)
    )  # パディングをピクセルに変換

    # タイトルエリアの高さ
    title_height = Pt(style.axis_label_pt * 2.5) if data.title else 0

    # 作業領域の計算
    work_left = container_left + padding_px
    work_top = container_top + title_height + padding_px
    work_width = container_width - padding_px * 2
    work_height = container_height - title_height - padding_px * 2

    # タイトル描画
    if data.title:
        _text_box(
            shapes,
            container_left,
            container_top,
            container_width,
            title_height,
            data.title,
            style.axis_label_pt + 2,
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    # 軸ラベル用のスペース
    x_axis_label_height = Pt(style.axis_label_pt * 2)
    y_axis_label_width = Pt(style.axis_label_pt * 4)

    # マトリクスの実際の描画領域
    matrix_left = work_left + y_axis_label_width
    matrix_top = work_top
    matrix_width = work_width - y_axis_label_width
    matrix_height = work_height - x_axis_label_height

    # 象限のサイズ
    quad_width = matrix_width // 2
    quad_height = matrix_height // 2

    # 象限の位置
    positions = {
        "bottom_left": (matrix_left, matrix_top + quad_height),
        "bottom_right": (matrix_left + quad_width, matrix_top + quad_height),
        "top_right": (matrix_left + quad_width, matrix_top),
        "top_left": (matrix_left, matrix_top),
    }

    # 象限の背景とコンテンツを描画
    quadrants = {
        "bottom_left": data.bottom_left,
        "bottom_right": data.bottom_right,
        "top_right": data.top_right,
        "top_left": data.top_left,
    }

    for quad_name, quad_data in quadrants.items():
        qx, qy = positions[quad_name]

        # 背景
        _draw_quadrant_bg(shapes, qx, qy, quad_width, quad_height, quad_data.color)

        # 象限内のパディング
        inner_padding = Pt(8)
        text_left = qx + inner_padding
        text_top = qy + inner_padding
        text_width = quad_width - inner_padding * 2
        text_height = quad_height - inner_padding * 2

        # タイトル
        title_h = Pt(style.quadrant_title_pt * 1.5)
        _text_box(
            shapes,
            text_left,
            text_top,
            text_width,
            title_h,
            quad_data.title,
            style.quadrant_title_pt,
            bold=True,
            align=PP_ALIGN.CENTER,
            v_anchor=MSO_ANCHOR.TOP,
        )

        # 説明（あれば）
        if quad_data.description:
            desc_top = text_top + title_h
            desc_height = text_height - title_h
            _text_box(
                shapes,
                text_left,
                desc_top,
                text_width,
                desc_height,
                quad_data.description,
                style.quadrant_desc_pt,
                align=PP_ALIGN.LEFT,
                v_anchor=MSO_ANCHOR.TOP,
            )

    # グリッド線（十字）
    center_x = matrix_left + quad_width
    center_y = matrix_top + quad_height

    # 縦線
    _draw_line(
        shapes,
        center_x,
        matrix_top,
        center_x,
        matrix_top + matrix_height,
        style.grid_color,
        style.grid_width_pt,
    )

    # 横線
    _draw_line(
        shapes,
        matrix_left,
        center_y,
        matrix_left + matrix_width,
        center_y,
        style.grid_color,
        style.grid_width_pt,
    )

    # 軸の描画
    axis_extend = Pt(10)  # 軸の延長

    # X軸（下）
    x_axis_y = matrix_top + matrix_height
    _draw_line(
        shapes,
        matrix_left - axis_extend,
        x_axis_y,
        matrix_left + matrix_width + axis_extend,
        x_axis_y,
        style.axis_color,
        style.axis_width_pt,
    )

    # Y軸（左）
    _draw_line(
        shapes,
        matrix_left,
        matrix_top - axis_extend,
        matrix_left,
        matrix_top + matrix_height + axis_extend,
        style.axis_color,
        style.axis_width_pt,
    )

    # 軸ラベル
    # X軸ラベル（中央下）
    _text_box(
        shapes,
        matrix_left,
        work_top + work_height - x_axis_label_height,
        matrix_width,
        x_axis_label_height,
        data.x_axis.label,
        style.axis_label_pt,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # X軸端ラベル
    end_label_width = Pt(40)
    x_label_offset = Pt(8)  # X軸ラベルを下に少しずらす
    # 左端（低）
    _text_box(
        shapes,
        matrix_left - end_label_width,
        x_axis_y + x_label_offset,
        end_label_width,
        Pt(style.axis_end_label_pt * 2),
        data.x_axis.low_label,
        style.axis_end_label_pt,
        align=PP_ALIGN.RIGHT,
    )
    # 右端（高）
    _text_box(
        shapes,
        matrix_left + matrix_width,
        x_axis_y + x_label_offset,
        end_label_width,
        Pt(style.axis_end_label_pt * 2),
        data.x_axis.high_label,
        style.axis_end_label_pt,
        align=PP_ALIGN.LEFT,
    )

    # Y軸ラベル（左中央、縦書き風に配置）
    _text_box(
        shapes,
        work_left,
        matrix_top,
        y_axis_label_width - Pt(5),
        matrix_height,
        data.y_axis.label,
        style.axis_label_pt,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # Y軸端ラベル
    end_label_height = Pt(style.axis_end_label_pt * 2)
    # 下端（低）
    _text_box(
        shapes,
        work_left,
        matrix_top + matrix_height - end_label_height,
        y_axis_label_width - Pt(5),
        end_label_height,
        data.y_axis.low_label,
        style.axis_end_label_pt,
        align=PP_ALIGN.CENTER,
        v_anchor=MSO_ANCHOR.TOP,
    )
    # 上端（高）
    _text_box(
        shapes,
        work_left,
        matrix_top,
        y_axis_label_width - Pt(5),
        end_label_height,
        data.y_axis.high_label,
        style.axis_end_label_pt,
        align=PP_ALIGN.CENTER,
        v_anchor=MSO_ANCHOR.BOTTOM,
    )

    # プロットポイントを描画
    if data.plot_points:
        point_radius = Pt(style.point_size_pt / 2)
        for point in data.plot_points:
            # 座標を計算（0.0〜1.0 → EMU）
            px = matrix_left + int(matrix_width * point.x)
            py = (
                matrix_top + matrix_height - int(matrix_height * point.y)
            )  # Y軸は上が高

            # 円を描画
            _draw_circle(shapes, px, py, point_radius, point.color)

            # ラベルを描画
            label_width = Pt(80)
            label_height = Pt(20)
            _text_box(
                shapes,
                px + point_radius,
                py - label_height // 2,
                label_width,
                label_height,
                point.label,
                style.axis_end_label_pt,
                align=PP_ALIGN.LEFT,
                color=point.color,
            )
