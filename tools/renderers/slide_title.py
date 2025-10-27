from __future__ import annotations

from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Pt
from pptx.dml.color import RGBColor

from ..utils import hex_to_rgb
from ..schemas.slide_title import SlideStyleSchema


def render(slide, data: SlideStyleSchema, geom: dict, context: dict):
    """
    左寄せタイトル。サブタイトルがあれば「 - 」で続け、やや小さい（既定 80%）フォントで同一行表示。

    geom: EMU dict {left, top, width, height}
    data: SlideStyleSchema with validated data
    """
    shapes = slide.shapes

    left, top, width, height = geom["left"], geom["top"], geom["width"], geom["height"]

    title = data.title.strip()
    subtitle = data.subtitle.strip() if data.subtitle else ""
    style = data.style or {}

    # --- スタイル既定 ---
    title_pt = style.title_pt
    if style.subtitle_pt is not None:
        sub_pt = max(8, style.subtitle_pt)
    else:
        sub_pt = max(8, int(round(title_pt * style.subtitle_scale)))

    # テーマ既定色の利用（任意）
    theme = (context or {}).get("theme") or {}
    color_hex = style.color if style.color != "#000000" else theme.get("font_color", "#000000")
    r, g, b = hex_to_rgb(color_hex)

    # --- テキストボックス ---
    tb = shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    try:
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
    except Exception:
        pass
    tf.clear()

    # 1段落内でメインとサブを runs で連結（同一行）
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.level = 0

    # メインタイトル
    run_main = p.add_run()
    run_main.text = title
    run_main.font.size = Pt(title_pt)
    run_main.font.bold = True
    run_main.font.color.rgb = RGBColor(r, g, b)

    # サブタイトル（" - " + subtitle）
    if subtitle:
        run_sep = p.add_run()
        run_sep.text = " - "
        run_sep.font.size = Pt(sub_pt)
        run_sep.font.bold = False
        run_sep.font.color.rgb = RGBColor(r, g, b)

        run_sub = p.add_run()
        run_sub.text = subtitle
        run_sub.font.size = Pt(sub_pt)
        run_sub.font.bold = False
        run_sub.font.color.rgb = RGBColor(r, g, b)

    # 区切り線（任意）
    divider = style.divider
    if divider and divider.on:
        from pptx.util import Emu  # 未使用でも import 可能（互換のため残置）
        dh_pct = divider.height_pct
        mt_pt = divider.margin_top_pt
        line_h = max(int(height * (dh_pct / 100.0)), int(Pt(1)))
        dy = top + height + int(Pt(mt_pt))
        dr, dg, db = hex_to_rgb(divider.color)

        rect = shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            left, dy, width, line_h
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(dr, dg, db)
        rect.line.fill.background()  # 枠線なし