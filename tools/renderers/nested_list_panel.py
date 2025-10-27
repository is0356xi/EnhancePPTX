"""
nested_list_panel ツール
- トップレベル: 角丸チップ（薄色塗り、オプションでアイコン）で見出し化
- ネスト(子):   枠線なしのテキスト行（行頭 "–"）をインデント表示
- YAML は「何を描くか（タイトル/項目/階層）」だけを渡し、デザインはツール側既定で吸収
- utils.py 仕様:
  - hex_to_rgb(hex: str) -> (r, g, b)
  - get_contrast_color(hex_bg_color: str) -> pptx.dml.color.RGBColor
"""
import math
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches

from ..utils import hex_to_rgb, get_contrast_color, normalize_color, rgb_to_hex, get_theme_color, mix_rgb_with_white
from ..schemas.nested_list_panel import NestedListPanelSchema


# =========================================================
# Helpers
# =========================================================
def _get_shapes_target(slide, context):
    return (context or {}).get("shapes_target", slide.shapes)


def _estimate_lines(text, width_emus, font_pt):
    """
    粗い行数見積: 幅(inch) と フォントサイズ(pt) から 1行収容文字数を近似し、行数を返す
    （AutoSizeを使わずにオーバーラップ回避するための近似）
    """
    if not text:
        return 0
    width_inch = max(0.01, width_emus / 914400.0)
    # 係数は経験値：1行あたり ≒ (幅 * 144) / pt
    chars_per_line = max(8.0, (width_inch * 144.0) / max(8.0, float(font_pt)))
    return max(1, int(math.ceil(len(str(text)) / chars_per_line)))


# =========================================================
# Render
# =========================================================
def render(slide, data: NestedListPanelSchema, geom: dict, context: dict):
    """
    data スキーマ（最小）:
      title: str (optional)     # パネルの見出し（任意）
      items:                     # 階層化された項目
        - text: str
          color: "#RRGGBB" (optional)   # 親チップの基調色（薄色に調整）
          icon: "DOT|TRIANGLE|NONE"     # 親チップの左アイコン（任意）
          children:
            - text: str
            - ...
      # 任意のヒント（必要なときだけ指定）
      # parent_pt: 14            # 親チップのフォントサイズ（pt）
      # child_pt:  12            # 子テキストのフォントサイズ（pt）
      # title_pt:  16            # タイトルのフォントサイズ（pt）
      # indent_em: 0.3           # 子テキストのインデント（inch）
      # icon_size_pt: 12         # アイコンサイズ（親チップ内）
      # gap_item_pt: 8           # アイテム間の上下ギャップ（pt）

    geom: { left, top, width, height } EMU
    """
    theme = (context or {}).get("theme", {}) or {}
    logger = (context or {}).get("logger", None)

    # --- data normalize ---
    title = data.title
    items = data.items

    title_pt   = data.title_pt
    parent_pt  = data.parent_pt
    child_pt   = data.child_pt
    indent_in  = data.indent_em
    icon_pt    = data.icon_size_pt
    gap_item_pt = data.gap_item_pt

    # colors
    font_rgb = normalize_color(theme.get("font_color", "#000000"))
    default_parent_rgb = get_theme_color(theme, "primary", "#0D6EFD")

    # geometry & paddings
    left, top, width, height = geom["left"], geom["top"], geom["width"], geom["height"]
    shapes = _get_shapes_target(slide, context)

    pad_x = max(int(width * 0.04), int(Pt(4)))
    pad_y = max(int(height * 0.06), int(Pt(4)))

    x0 = left + pad_x
    y  = top + pad_y
    w_inner = max(0, width - 2 * pad_x)

    # --- optional title ---
    if title:
        tb = shapes.add_textbox(x0, y, w_inner, int(Pt(title_pt * 1.6)))
        tf = tb.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT

        run = p.add_run()
        run.text = str(title)
        run.font.color.rgb = RGBColor(*font_rgb)
        try:
            run.font.size = Pt(title_pt)
        except Exception:
            pass

        y += int(Pt(title_pt * 1.6))  # タイトル分の高さ
        y += int(Pt(4))               # タイトルと1つ目のアイテムの間の余白

    # --- per item layout ---
    item_gap = int(Pt(gap_item_pt))
    chip_h   = int(Pt(max(18, parent_pt * 1.4)))  # 親チップの高さ（固定気味）
    chip_pad_x = int(Pt(6))
    chip_pad_y = int(Pt(2))
    icon_box_w = int(Pt(icon_pt + 6))  # アイコン領域の幅

    indent = int(Inches(indent_in))     # 子テキストのインデント（EMU）

    for idx, it in enumerate(items, start=1):
        # 溢れ回避（次のアイテム分の最低領域をざっくりチェック）
        if y + chip_h > top + height:
            if logger:
                logger.warning("nested_list_panel: 高さ不足のため一部項目を省略しました。")
            break
        text = it.text or ""
        children = it.children or []
        icon = (it.icon or "DOT").upper()
        base_rgb = normalize_color(it.color or default_parent_rgb)
        chip_rgb = mix_rgb_with_white(base_rgb, 0.88)  # 薄色化

        # === 親チップ（角丸＋塗り） ===
        chip = shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            x0, y, w_inner, chip_h
        )
        chip.fill.solid()
        chip.fill.fore_color.rgb = RGBColor(*chip_rgb)
        chip.line.fill.background()  # 罫線なし

        # アイコン（任意）
        icon_left = x0 + chip_pad_x
        icon_top  = y + int((chip_h - Pt(icon_pt)) / 2)
        text_left_in_chip = x0 + chip_pad_x
        if icon in ("DOT", "TRIANGLE"):
            if icon == "DOT":
                ico = shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.OVAL,
                    icon_left, icon_top, int(Pt(icon_pt)), int(Pt(icon_pt))
                )
            else:  # TRIANGLE
                ico = shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
                    icon_left, icon_top, int(Pt(icon_pt)), int(Pt(icon_pt))
                )
            # アイコン色は基調色（濃い目）
            ico.fill.solid()
            ico.fill.fore_color.rgb = RGBColor(*base_rgb)
            ico.line.fill.background()
            text_left_in_chip = x0 + icon_box_w  # テキスト開始位置をアイコン分右へ

        # 親テキスト（チップ内）
        tbp = shapes.add_textbox(
            text_left_in_chip,
            y + chip_pad_y,
            max(0, w_inner - (text_left_in_chip - x0) - chip_pad_x),
            max(0, chip_h - 2 * chip_pad_y)
        )
        tfp = tbp.text_frame
        tfp.clear()
        tfp.word_wrap = True
        tfp.vertical_anchor = MSO_ANCHOR.MIDDLE

        pp = tfp.paragraphs[0]
        pp.alignment = PP_ALIGN.LEFT
        runp = pp.add_run()
        runp.text = text

        # チップ背景に対するコントラスト色（utils は HEX→RGBColor を返す）
        chip_hex = rgb_to_hex(*chip_rgb)
        fg_color = get_contrast_color(chip_hex)
        runp.font.color.rgb = fg_color
        try:
            runp.font.size = Pt(parent_pt)
        except Exception:
            pass
        # 次行起点
        y += chip_h

        # === 子テキスト群（枠線なし／インデントのみ） ===
        if children:
            # 子全体の高さを見積もって確保（折返し込み）
            # 行高は child_pt * 1.3 を想定
            lines_total = 0
            for ch in children:
                t = ch.text or ""
                usable_w = max(1, w_inner - indent)  # インデント分を除いた幅
                lines_total += _estimate_lines(t, usable_w, child_pt)

            child_line_h = int(Pt(child_pt * 1.3))
            child_block_h = int(lines_total * child_line_h)

            # 高さ不足なら丸めて途中まで
            if y + child_block_h > top + height:
                # どれだけの行数が描けるか再計算
                avail_lines = max(0, (top + height - y) // child_line_h)
                # 描画可能な文だけ切り出し（単純に上から詰める）
                lines_written = 0
                child_textbox_h = int(avail_lines * child_line_h)
            else:
                avail_lines = lines_total
                lines_written = 0
                child_textbox_h = child_block_h

            if avail_lines > 0 and child_textbox_h > 0:
                tbc = shapes.add_textbox(x0 + indent, y + int(Pt(4)), max(0, w_inner - indent), child_textbox_h)
                tfc = tbc.text_frame
                tfc.clear()
                tfc.word_wrap = True
                tfc.vertical_anchor = MSO_ANCHOR.TOP

                # 1行目
                if children:
                    p0 = tfc.paragraphs[0]
                    p0.alignment = PP_ALIGN.LEFT
                    # 挿入関数
                    def _add_line(p, text_line):
                        run = p.add_run()
                        run.text = f"– {text_line}"   # 行頭記号（枠線なし）
                        run.font.color.rgb = RGBColor(*font_rgb)
                        try:
                            run.font.size = Pt(child_pt)
                        except Exception:
                            pass

                    # 子を順に追加（行数上限を超えたら打ち切り）
                    first = True
                    for ch in children:
                        t = ch.text or ""
                        # 行数見積（この子だけ）
                        usable_w = max(1, w_inner - indent)
                        need_lines = _estimate_lines(t, usable_w, child_pt)
                        if lines_written + need_lines > avail_lines:
                            # 余裕がなければ切断（簡略化：一行だけでも入るなら入れる）
                            if lines_written < avail_lines:
                                # 1行だけ入れる（粗い切り詰め）
                                t = t[: max(0, int(len(t) * (1.0 * (avail_lines - lines_written) / max(1, need_lines))))] + "…"
                                if first:
                                    _add_line(p0, t)
                                else:
                                    p = tfc.add_paragraph()
                                    p.alignment = PP_ALIGN.LEFT
                                    _add_line(p, t)
                                lines_written = avail_lines
                            break
                        else:
                            if first:
                                _add_line(p0, t)
                                first = False
                            else:
                                p = tfc.add_paragraph()
                                p.alignment = PP_ALIGN.LEFT
                                _add_line(p, t)
                            lines_written += need_lines

                y += child_textbox_h

        # アイテム間ギャップ
        y += item_gap

        # パネルからはみ出す場合は終了
        if y >= top + height:
            if logger:
                logger.warning("nested_list_panel: 高さ不足で切り上げました。")
            break