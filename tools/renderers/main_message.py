from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN

from ..utils import hex_to_rgb, get_contrast_color
from ..schemas.main_message import MainMessageSchema


def render(slide, data: MainMessageSchema, geom: dict, context: dict):
    """強調バナーを描画する"""
    # 1. バナー本体（背景の図形）を描画
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        geom["left"],
        geom["top"],
        geom["width"],
        geom["height"],
    )

    # 色と線
    bg_color_hex = data.color or context.get("theme", {}).get("primary", "#0D6EFD")
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*hex_to_rgb(bg_color_hex))
    shape.line.fill.background()

    # 角丸
    shape.adjustments[0] = int(data.corner_radius * 50000)

    # 文字色を決定
    font_color = get_contrast_color(bg_color_hex)

    # 3. テキストの描画
    text_frame = shape.text_frame
    text_frame.clear()

    # 左マージンを固定値に設定
    text_frame.margin_left = int(geom["width"] * 0.05)

    text_frame.margin_right = int(geom["width"] * 0.05)
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.NONE

    # メッセージ
    p_body = text_frame.paragraphs[0]
    run_body = p_body.add_run()
    run_body.text = data.main_message
    font_body = run_body.font
    font_body.size = Pt(data.style.font_size_body if data.style else 22)
    font_body.bold = True
    font_body.color.rgb = font_color

    # テキストの水平方向の配置
    align_str = (data.style.align if data.style else "left").upper()
    if hasattr(PP_ALIGN, align_str):
        p_body.alignment = getattr(PP_ALIGN, align_str)
