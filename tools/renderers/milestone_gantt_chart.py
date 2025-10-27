# -*- coding: utf-8 -*-
import colorsys
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Pt

from ..utils import hex_to_rgb, get_contrast_color, rgb_to_hex, mix_with_white
from ..schemas.milestone_gantt_chart import MilestoneGanttChartSchema

# 既定パレット（淡色/パステル系で6色）
PALETTE_HEX = ["#A7C8FF", "#AEE6C8", "#FFE7A3", "#F4B7B7", "#D7C2F2", "#FFC8A3"]

# ---- パステル生成ヘルパ ----


def _hsl_to_hex(h, s, l):
    r, g, b = colorsys.hls_to_rgb(h % 1.0, max(0.0, min(1.0, l)), max(0.0, min(1.0, s)))
    return rgb_to_hex(r * 255, g * 255, b * 255)


PASTEL_SEED_HEX = [
    "#FADADD",
    "#DDEEFF",
    "#FFF4CC",
    "#D8F2E6",
    "#EBDCFB",
    "#FFE6D5",
    "#D6F5F5",
    "#E7F0FF",
    "#FBE4E4",
    "#E9F9D2",
    "#E2E8F0",
    "#FDECC8",
]
_PASTEL_HUES = [
    0 / 360,
    15 / 360,
    30 / 360,
    45 / 360,
    55 / 360,
    80 / 360,
    120 / 360,
    150 / 360,
    180 / 360,
    200 / 360,
    220 / 360,
    240 / 360,
    260 / 360,
    285 / 360,
    310 / 360,
    335 / 360,
]
_PASTEL_VARIANTS = [(0.28, 0.90), (0.24, 0.92), (0.20, 0.94)]


def _pastel_by_index(i, s=None, l=None):
    if s is None or l is None:
        s_, l_ = _PASTEL_VARIANTS[(i // len(_PASTEL_HUES)) % len(_PASTEL_VARIANTS)]
        s = s if s is not None else s_
        l = l if l is not None else l_
    h = _PASTEL_HUES[i % len(_PASTEL_HUES)]
    return _hsl_to_hex(h, s, l)


def _shapes(slide, context):
    return (context or {}).get("shapes_target", slide.shapes)


# -------------------------
#  結合セル安全ユーティリティ
# -------------------------
def _tcPr(cell):
    return getattr(cell._tc, "tcPr", None)


def _is_spanned(cell) -> bool:
    """
    結合範囲内（左上以外）のダミーセルかどうか。
    PowerPoint 上は hMerge/vMerge のいずれかが付く。
    """
    tcPr = _tcPr(cell)
    if tcPr is None:
        return False
    hM = getattr(tcPr, "hMerge", None)
    vM = getattr(tcPr, "vMerge", None)
    return bool(hM) or bool(vM)


def _is_merge_origin(cell) -> bool:
    """
    結合の起点セル（左上・実体）かどうか。
    gridSpan/rowSpan が 1 より大きければ起点とみなせる。
    """
    tcPr = _tcPr(cell)
    if tcPr is None:
        return False
    gS = getattr(tcPr, "gridSpan", 1) or 1
    rS = getattr(tcPr, "rowSpan", 1) or 1
    return (gS > 1) or (rS > 1)


def _has_txBody(cell) -> bool:
    """
    セルに txBody（テキストボディ）が付いているか。spanned セルに付いていると危険。
    """
    return getattr(cell._tc, "txBody", None) is not None


def debug_dump_spanned_txbody(table):
    """
    spanned（ダミー）セルなのに txBody が付いているセルの座標を列挙する。
    """
    issues = []
    n_rows = len(table.rows)
    n_cols = len(table.columns)
    for r in range(n_rows):
        for c in range(n_cols):
            cell = table.cell(r, c)
            if _is_spanned(cell) and _has_txBody(cell):
                issues.append((r, c))
    return issues


def set_cell_text_safe(
    cell,
    text,
    *,
    font_size=10,
    bold=False,
    align=PP_ALIGN.LEFT,
    v_anchor=MSO_ANCHOR.MIDDLE,
    color_rgb=None,
    clear=True,
    left_indent_pt=None,
    word_wrap=None,
):
    """
    ダミー（spanned）セルには一切触らず、merge 後の実セル（または単独セル）にのみテキストを書き込む安全ラッパ。
    """
    if cell is None:
        raise ValueError(
            "set_cell_text_safe: cell が None です。merge() の戻り値が None の可能性があります。"
        )
    if _is_spanned(cell):
        raise RuntimeError(
            "spanned cell にテキストを設定しようとしました。merge 起点セル（左上）にのみ書いてください。"
        )

    tf = cell.text_frame
    if clear:
        tf.clear()

    try:
        tf.vertical_anchor = v_anchor
    except Exception:
        pass
    if word_wrap is not None:
        try:
            tf.word_wrap = bool(word_wrap)
        except Exception:
            pass

    p = tf.paragraphs[0]
    p.text = "" if text is None else str(text)
    p.font.size = Pt(font_size)
    p.font.bold = bool(bold)
    p.alignment = align
    if color_rgb is not None:
        p.font.color.rgb = color_rgb
    if left_indent_pt is not None:
        try:
            p.paragraph_format.left_indent = Pt(left_indent_pt)
        except Exception:
            pass


def _merge_ret(cell, other_cell):
    """
    python-pptx のバージョン差異対策。
    merge() の戻り値が None でも、左上セル（起点セル）を返す。
    """
    try:
        rv = cell.merge(other_cell)
        return rv if rv is not None else cell
    except Exception:
        # 例外時は最低限、起点セルを返す（結合できなかった場合）
        return cell


# -------------------------
#  データ色マッピング
# -------------------------
def _owner_hex_map(data: MilestoneGanttChartSchema):
    opt = data.options
    lighten_ratio = 0.70  # Default value since not in schema
    pastel_s = None
    pastel_l = None

    cmap = {}
    owners = data.owners or []
    for i, o in enumerate(owners):
        name = str(o.name or "").strip()
        col = o.color
        if not name:
            continue
        if isinstance(col, str) and col:
            if not col.startswith("#"):
                col = "#" + col
            try:
                pastel = mix_with_white(col, lighten_ratio)
            except Exception:
                pastel = _pastel_by_index(i, pastel_s, pastel_l)
        else:
            pastel = (
                PASTEL_SEED_HEX[i]
                if i < len(PASTEL_SEED_HEX)
                else _pastel_by_index(i, pastel_s, pastel_l)
            )
        cmap[name] = pastel

    idx = len(cmap)
    phases = data.phases or []
    for ph in phases:
        for t in ph.tasks or []:
            name = str(t.owner or "").strip() if hasattr(t, "owner") else ""
            if name and name not in cmap:
                pastel = (
                    PASTEL_SEED_HEX[idx]
                    if idx < len(PASTEL_SEED_HEX)
                    else _pastel_by_index(idx, pastel_s, pastel_l)
                )
                cmap[name] = pastel
                idx += 1
    return cmap


def _label_to_index(time_axis, key):
    # If key is numeric (int or float), return it as-is for position calculation
    if isinstance(key, (int, float)):
        return float(key)
    try:
        return float(time_axis.index(str(key)))
    except ValueError:
        return 0.0


# ---- 凡例用の小ヘルパ ----
def _legend_items(data: MilestoneGanttChartSchema, owner_hex, include_task_owners=True):
    items = []
    used = set()
    for o in data.owners or []:
        name = str(o.name or "").strip()
        if name and name in owner_hex and name not in used:
            items.append((name, owner_hex[name]))
            used.add(name)
    if include_task_owners:
        for ph in data.phases or []:
            for t in ph.tasks or []:
                name = str(t.owner or "").strip() if hasattr(t, "owner") else ""
                if name and name in owner_hex and name not in used:
                    items.append((name, owner_hex[name]))
                    used.add(name)
    return items


def _draw_owner_legend(
    slide, x, y, max_w, items, chip_pt=10, gap_chip_text_pt=6, gap_item_pt=14, font_pt=9
):
    """テキストの下にシンプルな凡例（色チップ＋名前）を横並びで描画。"""
    chip = int(Pt(chip_pt))
    gap_chip_text = int(Pt(gap_chip_text_pt))
    gap_item = int(Pt(gap_item_pt))
    font_size = Pt(font_pt)

    cursor_x = x
    cursor_y = y
    line_h = max(chip, int(Pt(font_pt + 2)))

    for name, hexcol in items:
        estimated_w = (
            chip + gap_chip_text + int(Pt(len(name) * (font_pt * 0.6))) + gap_item
        )
        if cursor_x + estimated_w > x + max_w:
            cursor_x = x
            cursor_y += line_h + int(Pt(4))

        rect = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, cursor_x, cursor_y, chip, chip
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(*hex_to_rgb(hexcol))
        rect.line.fill.background()

        tx = slide.shapes.add_textbox(
            cursor_x + chip + gap_chip_text, cursor_y - int(Pt(1)), int(Pt(240)), line_h
        )
        tf = tx.text_frame
        tf.clear()
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = font_size
        p.alignment = PP_ALIGN.LEFT

        cursor_x += estimated_w


# -------------------------
#  メイン描画
# -------------------------
def render(slide, data: MilestoneGanttChartSchema, geom: dict, context: dict):
    shapes_target = _shapes(slide, context)
    left, top, width, height = geom["left"], geom["top"], geom["width"], geom["height"]

    time_axis = data.time_axis or []
    milestones = data.milestones or []
    # Handle both phases and task_groups naming
    phases = data.phases or []
    if not phases and data.task_groups:
        # Convert task_groups to phases format for backward compatibility
        phases = [
            {"name": tg.group_name, "tasks": [task.model_dump() for task in tg.tasks]}
            for tg in data.task_groups
        ]
    opt = data.options

    if not time_axis:
        return

    label_col_pct = opt.label_col_pct
    header_fill_hex = opt.header_fill
    show_time_labels = opt.show_time_labels
    row_height_pt = opt.row_height_pt
    cell_shade_mode = opt.cell_shade.lower()  # "owner"|"phase"|"none"

    milestone_row_hex = opt.milestone_row_fill

    legend_opt = opt.legend
    legend_show = legend_opt.show
    legend_pos = legend_opt.position  # "below"|"above"|"right"
    legend_inc_task = legend_opt.include_task_owners

    # 追加の安全オプション
    assert_no_spanned_txbody = opt.assert_no_spanned_txbody

    owner_hex = _owner_hex_map(data)

    rows = (
        1
        + (1 if show_time_labels else 0)
        + sum(max(1, len(p.tasks or [])) for p in phases)
    )
    cols = 1 + len(time_axis)

    def _add_table_compat(shapes):
        if hasattr(shapes, "add_table"):
            return shapes.add_table(rows, cols, left, top, width, height)
        return slide.shapes.add_table(rows, cols, left, top, width, height)

    tbl_shape = _add_table_compat(shapes_target)
    table = tbl_shape.table

    # 行高
    for r in range(rows):
        table.rows[r].height = Pt(row_height_pt)

    # 列幅
    label_col_w = int(width * (label_col_pct / 100.0))
    table.columns[0].width = label_col_w
    if len(time_axis) > 0:
        slot_w = (width - label_col_w) / float(len(time_axis))
        for c in range(1, cols):
            table.columns[c].width = int(slot_w)
    else:
        slot_w = 0

    hdr_fill = RGBColor(*hex_to_rgb(header_fill_hex))
    milestone_fill = RGBColor(*hex_to_rgb(milestone_row_hex))

    # 全セル背景クリア（塗りは明示的に上書き）
    for r in range(rows):
        for c in range(cols):
            try:
                table.cell(r, c).fill.background()
            except Exception:
                pass

    # マイルストーン行の塗り（結合前に塗る）
    for c in range(cols):
        try:
            cell = table.cell(0, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = milestone_fill
        except Exception:
            pass
    # ---- ヘッダー（マイルストーン）: merge の戻り値に頼らず「起点セル」に書く ----
    try:
        header_cell = table.cell(0, 0)
        if cols > 1:
            header_cell = _merge_ret(header_cell, table.cell(0, cols - 1))
        set_cell_text_safe(
            header_cell,
            "マイルストーン",
            font_size=11,
            bold=True,
            align=PP_ALIGN.LEFT,
            v_anchor=MSO_ANCHOR.MIDDLE,
            left_indent_pt=6,
        )
    except Exception:
        pass

    current_row = 1

    # ---- 時間ラベル行（非結合セルのみ）----
    if show_time_labels:
        # 左端セル（ラベル列）は塗りのみ（テキスト不要）
        try:
            cell0 = table.cell(current_row, 0)
            cell0.fill.solid()
            cell0.fill.fore_color.rgb = hdr_fill
        except Exception:
            pass

        for i, label in enumerate(time_axis):
            cell = table.cell(current_row, 1 + i)
            try:
                cell.fill.solid()
                cell.fill.fore_color.rgb = hdr_fill
            except Exception:
                pass
            set_cell_text_safe(
                cell,
                str(label),
                font_size=10,
                align=PP_ALIGN.CENTER,
                v_anchor=MSO_ANCHOR.MIDDLE,
            )

        current_row += 1

    zebra_a = RGBColor(250, 250, 250)
    zebra_b = RGBColor(242, 242, 242)

    # ---- フェーズとタスク ----
    phase_index = 0
    for ph in phases:
        tasks = ph.tasks or []
        n_rows = max(1, len(tasks))
        start_row = current_row
        end_row = current_row + n_rows - 1

        phase_name = str(ph.name or "").strip()

        # ゼブラ背景（結合前に塗る）
        block_bg = zebra_a if (phase_index % 2 == 0) else zebra_b
        for rr in range(start_row, end_row + 1):
            for cc in range(cols):
                # 列0は縦結合予定。結合後 spanned を触らないため、ここでは rr>start_row をスキップ
                if cc == 0 and rr > start_row and n_rows > 1:
                    continue
                try:
                    cell = table.cell(rr, cc)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = block_bg
                except Exception:
                    pass

        # フェーズ名セルの結合とテキスト設定（戻り値に頼らず起点セルへ）
        try:
            phase_cell = table.cell(start_row, 0)
            if n_rows > 1:
                phase_cell = _merge_ret(phase_cell, table.cell(end_row, 0))
            if phase_name:
                set_cell_text_safe(
                    phase_cell,
                    phase_name,
                    font_size=10,
                    bold=True,
                    align=PP_ALIGN.CENTER,
                    v_anchor=MSO_ANCHOR.MIDDLE,
                )
        except Exception:
            # merge 失敗時も、開始セルのみに書く（範囲 spanned には触らない）
            if phase_name:
                try:
                    single_cell = table.cell(start_row, 0)
                    set_cell_text_safe(
                        single_cell,
                        phase_name,
                        font_size=10,
                        bold=True,
                        align=PP_ALIGN.CENTER,
                        v_anchor=MSO_ANCHOR.MIDDLE,
                    )
                except Exception:
                    pass
        # タスク行の描画
        for i in range(n_rows):
            task = tasks[i] if i < len(tasks) else None
            if task:
                start = int(task.start if hasattr(task, "start") else 0)
                end = int(task.end if hasattr(task, "end") else start + 1)
            else:
                start = 0
                end = 1
            start = max(0, min(start, len(time_axis) - 1)) if time_axis else 0
            end = max(start + 1, min(end, len(time_axis))) if time_axis else start + 1

            fill_rgb = None
            fill_hex = None
            if cell_shade_mode == "owner":
                owner = (
                    str(task.owner or "").strip()
                    if task and hasattr(task, "owner")
                    else ""
                )
                fill_hex = owner_hex.get(owner, PALETTE_HEX[0])
                fill_rgb = RGBColor(*hex_to_rgb(fill_hex))
            elif cell_shade_mode == "phase":
                base = block_bg
                try:
                    r = max(0, min(255, (base.rgb >> 16) - 10))
                    g = max(0, min(255, ((base.rgb >> 8) & 0xFF) - 10))
                    b = max(0, min(255, (base.rgb & 0xFF) - 10))
                    fill_rgb = RGBColor(r, g, b)
                except Exception:
                    fill_rgb = block_bg  # フォールバック

            # タスク名（非結合セルにのみ書く）
            if len(time_axis) > 0:
                name_cell = table.cell(current_row, 1 + start)
                name_color = None
                if fill_hex:
                    try:
                        contrast_hex = get_contrast_color(fill_hex)
                        name_color = RGBColor(*hex_to_rgb(contrast_hex))
                    except Exception:
                        name_color = None
                set_cell_text_safe(
                    name_cell,
                    str(task.name or "") if task and hasattr(task, "name") else "",
                    font_size=10,
                    align=PP_ALIGN.LEFT,
                    v_anchor=MSO_ANCHOR.MIDDLE,
                    color_rgb=name_color,
                )

            # ガント塗り（非結合セルのみ）
            if fill_rgb and len(time_axis) > 0:
                for c in range(start, end):
                    try:
                        cell = table.cell(current_row, 1 + c)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = fill_rgb
                    except Exception:
                        pass
            current_row += 1

        phase_index += 1

    # ---- spanned セルに txBody が付いていないかの検査（任意）----
    if assert_no_spanned_txbody:
        issues = debug_dump_spanned_txbody(table)
        if issues:
            raise RuntimeError(
                f"spanned セルに txBody が付いています: {issues}. spanned へ .text_frame を触っていないか確認してください。"
            )

    # ---- マイルストーン（★）とラベル ----
    top_row_h = table.rows[0].height
    label_col_w = table.columns[0].width
    slot_w = table.columns[1].width if len(time_axis) > 0 else 0

    star_size = int(Pt(14))
    star_y = top + int(top_row_h * 0.55) - star_size // 2

    for ms in milestones:
        idx = _label_to_index(time_axis, ms.time_index)
        # Clamp the index to valid range, but preserve fractional values
        idx = max(0.0, min(idx, float(len(time_axis) - 1))) if time_axis else 0.0
        cx = int(left + label_col_w + idx * slot_w + slot_w / 2)

        star = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.STAR_5_POINT,
            cx - star_size // 2,
            star_y,
            star_size,
            star_size,
        )
        star.fill.solid()
        star.fill.fore_color.rgb = RGBColor(255, 193, 7)
        star.line.fill.background()

        box = slide.shapes.add_textbox(
            cx - int(Pt(60)), star_y - int(Pt(20)), int(Pt(120)), int(Pt(18))
        )
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = str(ms.label or "")
        p.font.size = Pt(9)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    # ---- 凡例 ----
    if legend_show:
        items = _legend_items(data, owner_hex, include_task_owners=legend_inc_task)
        if items:
            margin = int(Pt(6))
            if legend_pos == "above":
                legend_x = left
                legend_y = top - int(Pt(20))
            elif legend_pos == "right":
                legend_x = left + width + margin
                legend_y = top + int(Pt(2))
            else:  # "below"
                legend_x = left
                legend_y = top + height + margin
            _draw_owner_legend(slide, legend_x, legend_y, width, items)

    return
