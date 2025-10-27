from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml import parse_xml

from ..utils import hex_to_rgb
from ..schemas.component_diagram import Schema as ComponentsSchema

NODE_SHAPES = {
    "user": MSO_SHAPE.ROUND_1_RECTANGLE,
    "rect": MSO_SHAPE.RECTANGLE,
}


# ------------------------------------------------------------
# 位置計算ユーティリティ
# ------------------------------------------------------------
def _rel_to_abs(geom, rel):
    left, top, width, height = geom["left"], geom["top"], geom["width"], geom["height"]
    ax = left + int(width * (rel.x / 100.0))
    ay = top + int(height * (rel.y / 100.0))
    aw = int(width * (rel.w / 100.0))
    ah = int(height * (rel.h / 100.0))
    return {"x": ax, "y": ay, "w": aw, "h": ah}


def _calc_conn_points_and_sites(from_g, to_g, margin=0):
    """
    図形 from_g → to_g の位置関係を解析し、
    (start_xy, end_xy), (begin_site, end_site), connector_type を返す。
    site: 0=上,1=左,2=下,3=右
    """
    fx, fy, fw, fh = from_g["x"], from_g["y"], from_g["w"], from_g["h"]
    tx, ty, tw, th = to_g["x"], to_g["y"], to_g["w"], to_g["h"]

    f_left, f_right = fx, fx + fw
    f_top, f_bottom = fy, fy + fh
    t_left, t_right = tx, tx + tw
    t_top, t_bottom = ty, ty + th

    fc_x, fc_y = fx + fw / 2, fy + fh / 2
    tc_x, tc_y = tx + tw / 2, ty + th / 2

    dx = abs(tc_x - fc_x)
    dy = abs(tc_y - fc_y)

    # ほぼ水平かどうかの閾値（図形の高さの平均の半分より縦のズレが小さいか）
    is_horizontal_aligned = dy < (fh + th) / 4
    # ほぼ垂直かどうかの閾値（図形の幅の平均の半分より横のズレが小さいか）
    is_vertical_aligned = dx < (fw + tw) / 4

    # 1) ほぼ水平に並んでいて、かつ左右に分離している場合 → 直線
    if is_horizontal_aligned and (f_right <= t_left or t_right <= f_left):
        connector_type = MSO_CONNECTOR.STRAIGHT
        if f_right <= t_left:  # from -> to
            return (
                ((f_right + margin, fc_y), (t_left - margin, tc_y)),
                (3, 1),
                connector_type,
            )
        else:  # to -> from
            return (
                ((f_left - margin, fc_y), (t_right + margin, tc_y)),
                (1, 3),
                connector_type,
            )

    # 2) ほぼ垂直に並んでいて、かつ上下に分離している場合 → 直線
    if is_vertical_aligned and (f_bottom <= t_top or t_bottom <= f_top):
        connector_type = MSO_CONNECTOR.STRAIGHT
        if f_bottom <= t_top:  # from -> to
            return (
                ((fc_x, f_bottom + margin), (tc_x, t_top - margin)),
                (2, 0),
                connector_type,
            )
        else:  # to -> from
            return (
                ((fc_x, f_top - margin), (tc_x, t_bottom + margin)),
                (0, 2),
                connector_type,
            )

    # 3) 上記以外（斜め、または重なっている）→ 常にカギ型
    connector_type = MSO_CONNECTOR.ELBOW
    if dx >= dy:  # 横方向優勢
        if tc_x > fc_x:  # to が右側
            return ((fc_x, fc_y), (tc_x, tc_y)), (3, 1), connector_type
        else:  # to が左側
            return ((fc_x, fc_y), (tc_x, tc_y)), (1, 3), connector_type
    else:  # 縦方向優勢
        if tc_y > fc_y:  # to が下側
            return ((fc_x, fc_y), (tc_x, tc_y)), (2, 0), connector_type
        else:  # to が上側
            return ((fc_x, fc_y), (tc_x, tc_y)), (0, 2), connector_type


# ------------------------------------------------------------
# コネクタ生成（矢印向き修正版）
# ------------------------------------------------------------
def _add_connector(
    slide,
    start,
    end,
    style,
    connector_type,
    from_shape=None,
    to_shape=None,
    site_pair=None,
):
    conn = slide.shapes.add_connector(
        connector_type, int(start[0]), int(start[1]), int(end[0]), int(end[1])
    )

    if from_shape and to_shape and site_pair:
        try:
            conn.begin_connect(from_shape, site_pair[0])
            conn.end_connect(to_shape, site_pair[1])
        except Exception:
            pass

    line = conn.line
    line.width = Pt(style.pt)
    line.color.rgb = RGBColor(*hex_to_rgb(style.color))
    if style.dash == "dashed":
        line.dash_style = MSO_LINE_DASH_STYLE.DASH

    # Default to end arrow style for now
    ln = line._get_or_add_ln()
    ns = 'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
    ln.append(parse_xml(f'<a:tailEnd type="arrow" {ns}/>'))

    return conn


# ------------------------------------------------------------
# ラベル描画
# ------------------------------------------------------------
def _draw_label(slide, text, center_xy, context):
    w = Pt(100)
    h = Pt(28)
    left = center_xy[0] - int(w / 2)
    top = center_xy[1] - int(h / 2)
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(10)

    default_font_color = context.get("theme", {}).get("font_color", "#000000")
    r.font.color.rgb = RGBColor(*hex_to_rgb(default_font_color))

    return tb


# ------------------------------------------------------------
# メイン描画関数
# ------------------------------------------------------------
def render(slide, data: ComponentsSchema, geom: dict, context: dict):
    node_map = {}
    default_font_color = context.get("theme", {}).get("font_color", "#000000")

    # --- ノード ---
    for nd in data.nodes:
        kind = nd.kind
        if kind not in NODE_SHAPES:
            kind = "rect"

        g = _rel_to_abs(geom, nd.pos)
        shape = slide.shapes.add_shape(
            NODE_SHAPES[kind], g["x"], g["y"], g["w"], g["h"]
        )

        st = nd.style
        fill_hex = (
            st.fill or ("#F0F4FF" if kind == "user" else "#FFFFFF")
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*hex_to_rgb(fill_hex))
        if st.stroke:
            shape.line.color.rgb = RGBColor(*hex_to_rgb(st.stroke))
        shape.line.width = Pt(st.stroke_pt)

        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = nd.label or ""
        r.font.size = Pt(st.font_size or 11)

        font_color_hex = st.font_color or default_font_color
        r.font.color.rgb = RGBColor(*hex_to_rgb(font_color_hex))

        node_map[nd.id] = {"geom": g, "shape": shape}

    # --- コネクタ ---
    for c in data.connectors:
        fn = node_map.get(c.from_)
        tn = node_map.get(c.to)
        if not fn or not tn:
            continue

        (start, end), sites, connector_type = _calc_conn_points_and_sites(
            fn["geom"], tn["geom"]
        )
        conn = _add_connector(
            slide,
            start,
            end,
            c.style,
            connector_type,
            from_shape=fn["shape"],
            to_shape=tn["shape"],
            site_pair=sites,
        )

        if c.label:
            mid = ((start[0] + end[0]) / 2, (start[1] + end[1]) / 2)
            _draw_label(slide, c.label, mid, context)