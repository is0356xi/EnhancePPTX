import math
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Pt

from ..utils import hex_to_rgb, get_contrast_color, normalize_color, get_theme_color, rgb_to_hex
from ..schemas.progress_bar import ProgressBarSchema


def _choose_bar_color(theme: dict, pct: float, target: float):
    """進捗値に応じてバー色を決定"""
    success = get_theme_color(theme, "success", "#198754")  # green
    primary = get_theme_color(theme, "primary", "#0D6EFD")  # blue
    warning = get_theme_color(theme, "warning", "#FFC107")  # amber
    danger  = get_theme_color(theme, "danger",  "#DC3545")  # red

    if pct >= target:
        return success
    if pct >= 70:
        return primary
    if pct >= 40:
        return warning
    return danger


def _get_shapes_target(slide, context):
    return (context or {}).get("shapes_target", slide.shapes)


def render(slide, data: ProgressBarSchema, geom: dict, context: dict):
    """Progress bar renderer with Pydantic validation"""
    theme = (context or {}).get("theme", {}) or {}

    # Data is already validated by Pydantic
    title = data.title
    text_suffix = data.text
    title_pt = data.title_pt
    suffix_pt = data.suffix_pt
    current_pct = max(0.0, min(100.0, data.current_pct))  # clamp to 0-100
    target_pct = data.target_pct

    # --- layout within geom ---
    left, top, width, height = geom["left"], geom["top"], geom["width"], geom["height"]
    shapes = _get_shapes_target(slide, context)

    # ・左右/上下のパディング（比率 + 絶対最小）
    pad_x = max(int(width * 0.05), int(Pt(4)))   # 最低 4pt
    pad_y = max(int(height * 0.08), int(Pt(4)))  # 最低 4pt

    # ・ギャップ（タイトル下の空き）
    gap_y = max(int(height * 0.06), int(Pt(4)))  # 最低 4pt

    # ===== Title area estimation =====
    width_inch = max(0.01, width / 914400.0)
    chars_per_line = max(8.0, (width_inch * 144.0) / max(8.0, float(title_pt)))

    # サフィックスは小さい文字なので、相対長をスケール
    suffix_scale = (float(suffix_pt) / float(title_pt)) if title_pt > 0 else 0.6
    logical_len = len(title)
    if text_suffix:
        logical_len += 3  # " - "
        logical_len += int(len(str(text_suffix)) * max(0.4, min(1.0, suffix_scale)))

    est_lines = max(1, int(math.ceil(logical_len / chars_per_line)))
    est_lines = min(est_lines, 4)  # 上限 4 行

    # 行高：title_pt をベースに 1.3 行送り
    line_height_emus = int(Pt(title_pt * 1.3))
    title_h = pad_y + est_lines * line_height_emus
    # タイトル領域が全高の 60% を超えないように制限
    title_h = min(title_h, int(height * 0.6))

    # バーの最低高さ（8pt 相当）
    min_bar_h = int(Pt(8))

    # 実際に使える高さから算出
    bar_h = height - title_h - gap_y
    if bar_h < min_bar_h:
        # タイトルが大きすぎるので縮めてでもバー最小高＋ギャップを確保
        title_h = max(0, height - gap_y - min_bar_h)
        bar_h = min_bar_h

    # --- Title (+ suffix text) ---
    if title or text_suffix:
        tb = shapes.add_textbox(left + pad_x, top + pad_y, max(0, width - 2*pad_x), max(0, title_h - pad_y))
        tf = tb.text_frame
        tf.clear()
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT

        font_rgb = normalize_color(theme.get("font_color", "#000000"))

        # メインタイトル
        if title:
            run_title = p.add_run()
            run_title.text = title
            run_title.font.color.rgb = RGBColor(*font_rgb)
            try:
                run_title.font.size = Pt(title_pt)
            except Exception:
                pass

        # suffix（小さめ）：" - " を挟んで同一行
        if text_suffix:
            run_sep = p.add_run()
            run_sep.text = " - "
            run_sep.font.color.rgb = RGBColor(*font_rgb)
            try:
                run_sep.font.size = Pt(title_pt)
            except Exception:
                pass

            run_suffix = p.add_run()
            run_suffix.text = str(text_suffix)
            run_suffix.font.color.rgb = RGBColor(*font_rgb)
            try:
                run_suffix.font.size = Pt(suffix_pt)
            except Exception:
                pass

    # --- Bar background (rounded) ---
    bar_top = top + title_h + gap_y
    bar_left = left + pad_x
    bar_width = max(0, width - 2*pad_x)

    bar_bg = shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        bar_left, bar_top, bar_width, bar_h
    )
    bg_rgb = normalize_color("#E9ECEF")  # 薄いグレー
    bar_bg.fill.solid()
    bar_bg.fill.fore_color.rgb = RGBColor(*bg_rgb)
    bar_bg.line.fill.background()

    # --- Bar fill (progress) ---
    progress_ratio = min(max(current_pct / target_pct, 0.0), 1.0)
    fill_width = int(bar_width * progress_ratio)

    bar_color = _choose_bar_color(theme, current_pct, target_pct)
    if fill_width > 0:
        bar_fill = shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            bar_left, bar_top, fill_width, bar_h
        )
        bar_fill.fill.solid()
        bar_fill.fill.fore_color.rgb = RGBColor(*bar_color)
        bar_fill.line.fill.background()

    # --- Percentage text over the bar ---
    pct_tb = shapes.add_textbox(bar_left, bar_top, bar_width, bar_h)
    pct_tf = pct_tb.text_frame
    pct_tf.clear()
    pct_tf.word_wrap = False
    pct_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = pct_tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    pct_run = p.add_run()
    pct_run.text = f"{round(current_pct)}%"

    # 背景とのコントラスト色
    base_rgb = bar_color if progress_ratio > 0.5 else bg_rgb
    base_hex = rgb_to_hex(*base_rgb)
    fg_color = get_contrast_color(base_hex)  # RGBColor
    pct_run.font.color.rgb = fg_color