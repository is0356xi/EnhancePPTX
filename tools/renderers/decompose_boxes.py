from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt
from pptx.dml.color import RGBColor

from ..schemas.decompose_boxes import DecomposeBoxesSchema


def render(slide, data: DecomposeBoxesSchema, geom: dict, context: dict):
    """
    決定論的レイアウト：
    - 列は最大深度+1。
    - 各列は等幅、列間ギャップは固定率。
    - 列ごとに上部ヘッダー（テキストボックス、背景なし）。
    - 親高さ = 子合計高さ（兄弟間の垂直ギャップは親にも子にも含めない）。
    - コネクタは描画しない。
    - 色は列0: 薄青 (#DDEBF7) / 列1以降: 薄灰 (#F2F2F2)、枠線なし。
    """
    shapes = context.get("shapes_target", slide.shapes)
    # Normalize root: allow DecomposeBoxNode or list[DecomposeBoxNode];
    # convert pydantic models to plain dicts for renderer convenience.
    raw_root = data.root
    if isinstance(raw_root, list):
        converted = []
        for r in raw_root:
            if hasattr(r, "model_dump"):
                converted.append(r.model_dump())
            else:
                converted.append(r)
        root_data = converted
    else:
        if hasattr(raw_root, "model_dump"):
            root_data = raw_root.model_dump()
        else:
            root_data = raw_root
    headers = data.column_headers

    # -------------------------
    # 既定パラメータ（必要に応じて utils へ移せます）
    # -------------------------
    col_gap_pct = 0.03  # 列間ギャップ（描画領域幅に対する比）
    row_gap_pct = 0.01  # 兄弟間ギャップ（描画領域高に対する比）
    header_band_pct = 0.08  # 列ヘッダー高さ（描画領域高に対する比）
    font_title = Pt(12)
    font_header = Pt(10)

    # 既定色（ご要望の通り）
    LIGHT_BLUE = RGBColor(0xDD, 0xEB, 0xF7)  # #DDEBF7
    LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)  # #F2F2F2

    # -------------------------
    # 深さ＆列数
    # -------------------------
    def depth(node):
        # node may be a dict-like node or a list of nodes (top-level root list)
        if isinstance(node, list):
            if not node:
                return 0
            return max(depth(n) for n in node)
        # guard: if node isn't a mapping-like object, treat as leaf
        if not isinstance(node, dict):
            return 0
        children = node.get("children")
        if not children:
            return 0
        return 1 + max(depth(ch) for ch in children)

    D = depth(root_data)
    cols = D + 1

    # -------------------------
    # 幾何（列幅・位置・コンテンツ帯）
    # -------------------------
    col_gap = int(geom["width"] * col_gap_pct) if cols > 1 else 0
    total_col_gaps = col_gap * (cols - 1)

    # 左側の列（col=0）をデフォルトで狭くする
    if cols > 1:
        # 左側の列を全体の20%、残りの列を等分
        first_col_ratio = 0.2
        remaining_width = geom["width"] - total_col_gaps
        first_col_w = int(remaining_width * first_col_ratio)
        other_col_w = (remaining_width - first_col_w) // (cols - 1)
        col_widths = [first_col_w] + [other_col_w] * (cols - 1)
    else:
        # 列が1つしかない場合は従来通り
        col_w = geom["width"]
        col_widths = [col_w]

    # Only reserve header band if explicit column_headers are provided
    header_h = int(geom["height"] * header_band_pct) if headers else 0
    content_top = geom["top"] + header_h
    content_h = geom["height"] - header_h

    def col_x(c):
        if c == 0:
            return geom["left"]
        else:
            x = geom["left"] + col_widths[0] + col_gap
            for i in range(1, c):
                x += col_widths[i] + col_gap
            return x

    def col_width(c):
        return col_widths[c] if c < len(col_widths) else col_widths[-1]

    # 兄弟間の縦ギャップ（固定）
    gap_y = int(geom["height"] * row_gap_pct)

    # -------------------------
    # 列ヘッダー（背景なしのテキスト）
    # -------------------------
    # Draw column headers only when the user provided them.
    # If the provided `column_headers` list is shorter than the number of
    # columns, do not auto-fill or draw placeholders for the missing headers.
    if headers:
        for c in range(cols):
            if c >= len(headers):
                # No header string supplied for this column — skip drawing.
                continue
            tb = shapes.add_textbox(col_x(c), geom["top"], col_width(c), header_h)
            tf = tb.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = headers[c]
            p.font.size = font_header
            p.font.bold = True

    # -------------------------
    # レイアウト：親→子に高さを分配（合計一致）
    # -------------------------
    boxes = []  # (node, col, x, y, w, h)

    def layout(node, col, x, y, w, h):
        # この node の箱（列 col のコンテンツ帯内）
        y = max(y, content_top)
        bottom = content_top + content_h
        h = min(h, bottom - y)
        boxes.append((node, col, x, y, w, h))

        children = node.get("children") or []
        if not children:
            return

        n = len(children)
        # 重み（省略時は等分）
        weights = [float(ch.get("weight", 1.0)) for ch in children]
        total_w = sum(wi for wi in weights if wi > 0) or float(n)
        # 子同士のギャップを除外した可用高さ
        usable = h - gap_y * (n - 1)
        if usable < 0:
            # ギャップが大きすぎる異常系：ギャップなしで強制配分
            usable = h
            local_gap = 0
        else:
            local_gap = gap_y

        cursor_y = y
        for i, ch in enumerate(children):
            wi = max(weights[i], 0.0)
            frac = (wi / total_w) if total_w > 0 else (1.0 / n)
            # 端数は最後の子に吸収して合計一致
            if i < n - 1:
                ch_h = int(round(usable * frac))
            else:
                ch_h = y + h - cursor_y
            child_col = min(col + 1, cols - 1)
            child_x = col_x(child_col)
            layout(ch, child_col, child_x, cursor_y, col_width(child_col), ch_h)
            cursor_y += ch_h + (local_gap if i < n - 1 else 0)

    # ルートを全高で配置
    # もし root_data がリストで渡されたら、左列 (col=0) に複数のトップレベル箱を分割して配置する
    # （各要素に 'weight' があれば比率配分に使う）
    if isinstance(root_data, list):
        n = len(root_data)
        if n > 0:
            weights = [float(item.get("weight", 1.0)) for item in root_data]
            total_w = sum(w for w in weights if w > 0) or float(n)
            # 子同士のギャップを除いた可用高さ
            usable = content_h - gap_y * (n - 1)
            if usable < 0:
                usable = content_h
                local_gap = 0
            else:
                local_gap = gap_y

            cursor_y = content_top
            for i, item in enumerate(root_data):
                wi = max(weights[i], 0.0)
                frac = (wi / total_w) if total_w > 0 else (1.0 / n)
                # 最後の要素は端数を吸収して合計一致させる
                if i < n - 1:
                    h_i = int(round(usable * frac))
                else:
                    h_i = content_top + content_h - cursor_y
                layout(item, 0, col_x(0), cursor_y, col_width(0), h_i)
                cursor_y += h_i + (local_gap if i < n - 1 else 0)
    else:
        layout(root_data, 0, col_x(0), content_top, col_width(0), content_h)

    # -------------------------
    # 描画：矩形＋テキスト（枠線なし）
    # -------------------------
    def add_box(node, col, x, y, w, h):
        shp = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shp.fill.solid()
        # 列 0 は薄青、それ以降は薄灰
        shp.fill.fore_color.rgb = LIGHT_BLUE if col == 0 else LIGHT_GRAY
        # 枠線なし
        shp.line.fill.background()
        # 角丸（やりすぎない程度）※環境により 0〜1 相当の調整値
        try:
            shp.adjustments[0] = 0.08
        except Exception:
            pass

        # テキスト（スキーマの 'name' を使用）
        tf = shp.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = node.get("name", "")
        p.font.size = font_title
        p.font.bold = True if col <= 1 else False
        tf.word_wrap = True

    for node, col, x, y, w, h in boxes:
        add_box(node, col, x, y, w, h)
