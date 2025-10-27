from __future__ import annotations

from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Pt

from ..utils import hex_to_rgb, get_contrast_color  # get_contrast_color: HEX -> RGBColor
from ..schemas.pyramid_hierarchy import PyramidHierarchySchema

# ====== 可読性・見た目の既定（必要ならここだけ弄ればOK） ======
EMU_PER_INCH = 914400
PT_PER_INCH = 72.0

# 連続した斜辺（ギャップをまたいでも一直線）を作る設定
RATIO_DEFAULT = 0.80     # r=0.8 の思想（Apex 初期値の算出に使用）
MIN_TOP_RATIO = 0.22     # 最上点（Apex）の最小幅: 全体幅の22%（読める下限）

SIDE_INSET = 0.03        # 親枠から左右3%内側に（外枠と干渉させない）
GAP_PCT = 2.0            # 段間ギャップ %
OVERLAP_PX_RATIO = 0.0015  # 段境界の隙間対策：横方向に重ねる比率（≒0.15%）

MIN_FONT_PT, MAX_FONT_PT = 12, 26  # ラベルの自動フォント pt の上下限

# ====== 内部ユーティリティ ======

def _sanitize_palette(palette):
    if not palette:
        return []
    out = []
    for c in palette:
        if not c:
            continue
        c = c.strip()
        if not c.startswith("#"):
            c = "#" + c
        out.append(c.upper())
    return out

def _lighten_hex(hx: str, factor: float = 0.15) -> str:
    hx = hx.lstrip("#")
    r = int(hx[0:2], 16); g = int(hx[2:4], 16); b = int(hx[4:6], 16)
    def blend(c): return int(round(c + (255 - c) * max(0.0, min(1.0, factor))))
    nr, ng, nb = blend(r), blend(g), blend(b)
    return "#{:02X}{:02X}{:02X}".format(nr, ng, nb)

def _auto_color(i: int, palette: list[str], theme: dict) -> str:
    if palette:
        return palette[i % len(palette)]
    base = (theme or {}).get("primary", "#0D6EFD")
    # 下段ほど濃く、上段ほどわずかに明るく
    return _lighten_hex(base, factor=0.12 * i)

def _auto_font_pt(h_emu: int, override_pt: int | None) -> int:
    """段の高さからフォントサイズを自動算出（2行前提）"""
    if override_pt:
        return int(override_pt)
    h_in = h_emu / EMU_PER_INCH
    pt = int(round(h_in * PT_PER_INCH * 0.45))
    return max(MIN_FONT_PT, min(MAX_FONT_PT, pt))

# ====== メイン描画 ======

def render(slide, data: PyramidHierarchySchema, geom: dict, context: dict):
    """
    ピラミッド（下→上）を描画。geom は EMU 完成形: {"left","top","width","height"}。
    YAML は抽象（title / levels[*] = {name, subtitle?, items?}）。
    形状の細部はツールが自動最適化し、側面は"1本の直線"に見えるよう連続化する。
    """
    shapes = slide.shapes
    theme = context.get("theme", {})
    logger = context.get("logger")

    # --- 入力 ---
    left, top, width, height = geom["left"], geom["top"], geom["width"], geom["height"]
    title = data.title
    levels = data.levels
    if not levels:
        if logger:
            logger.warning("pyramid_hierarchy: levels is empty; skip")
        return
    n = len(levels)

    # スタイル
    style = data.style
    palette = _sanitize_palette(style.palette)
    outline = style.outline
    font_style = style.font
    font_bold = font_style.bold
    font_align_key = font_style.align.upper()
    font_align = getattr(PP_ALIGN, font_align_key, PP_ALIGN.CENTER)
    font_override_pt = font_style.size_pt

    # --- 寸法（高さ・ギャップ・バンド） ---
    gap = int(height * (GAP_PCT / 100.0))
    total_gap = gap * (n - 1)
    h_eff = max(height - total_gap, 0)
    h_each = [h_eff // n] * n  # 均等配分（読みやすさ優先）

    # 各段の y 座標（下から上へ）
    y_cursor = top + height
    bands = []
    for i in range(n):
        h = h_each[i]
        y_bottom = y_cursor
        y_top = y_bottom - h
        bands.append((y_top, y_bottom))
        y_cursor = y_top - (gap if i < n - 1 else 0)

    # --- 連続傾斜の幅算出（ギャップ込みで直線補間） ---
    cx = left + width // 2
    inner_width = int(width * (1.0 - 2 * SIDE_INSET))  # 左右インセットで外枠回避

    # 連続直線の基準点（最上点↔最下点）
    y_apex = bands[-1][0]   # 最上段の上端
    y_base = bands[0][1]    # 最下段の下端
    H_total = max(1, y_base - y_apex)

    # r=0.8 の思想を Apex 初期値に反映しつつ、読みやすさの下限も確保
    w_base = inner_width
    w_apex_from_r = int(round(inner_width * (RATIO_DEFAULT ** n)))
    w_apex_min     = int(inner_width * MIN_TOP_RATIO)
    w_apex = max(w_apex_from_r, w_apex_min)

    def width_at_y(y: int) -> int:
        """高さ y における幅（線形補間）— ギャップをまたいでも連続"""
        t = (y - y_apex) / H_total
        return int(round(w_apex + (w_base - w_apex) * t))

    # 各段の上辺/下辺幅
    widths_top    = [width_at_y(bands[i][0]) for i in range(n)]
    widths_bottom = [width_at_y(bands[i][1]) for i in range(n)]

    # 丸め誤差の保険：厳密に「上辺 < 下辺」にする（最低 1px 差）
    min_step = max(int(width * 0.002), 1)
    for i in range(n):
        if widths_bottom[i] <= widths_top[i]:
            widths_bottom[i] = widths_top[i] + min_step

    # 段境界の"ヘアライン"対策：横方向にわずかに重ねる
    overlap_px = max(int(width * OVERLAP_PX_RATIO), 1)
    def _expand_pair(w_top_i: int, w_bottom_i: int) -> tuple[int, int]:
        return w_top_i + 2 * overlap_px, w_bottom_i + 2 * overlap_px

    # --- 描画（下→上） ---
    for i, lvl in enumerate(levels):
        # 幅（オーバーラップ適用）
        w_top_i, w_bottom_i = _expand_pair(widths_top[i], widths_bottom[i])

        # 台形の四隅（EMU、左右対称）
        y_top_i, y_bottom_i = bands[i]
        xlb = cx - w_bottom_i // 2; xrb = cx + w_bottom_i // 2
        xlt = cx - w_top_i    // 2; xrt = cx + w_top_i    // 2

        fb = shapes.build_freeform(xlb, y_bottom_i)
        fb.add_line_segments([
            (xrb, y_bottom_i),   # 底辺右端
            (xrt, y_top_i),      # 右上
            (xlt, y_top_i),      # 左上
            (xlb, y_bottom_i),   # 左下（close直前）
        ], close=True)
        shape = fb.convert_to_shape()

        # 塗り - using auto color since no color in schema level
        color_hex = _auto_color(i, palette, theme)
        rr, gg, bb = hex_to_rgb(color_hex)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(rr, gg, bb)

        # 境界線（白の細線で段差をくっきり）
        if outline:
            oc_hex = outline.color
            or_, og, ob = hex_to_rgb(oc_hex)
            shape.line.width = Pt(outline.width_pt)
            shape.line.color.rgb = RGBColor(or_, og, ob)

        # 帯内ラベル（name / subtitle）
        tf = shape.text_frame
        tf.word_wrap = True
        tf.clear()
        p = tf.paragraphs[0]
        name = lvl.name or ""
        subtitle = lvl.subtitle or ""
        p.text = f"{name}\n{subtitle}".strip()

        fnt_pt = _auto_font_pt(h_each[i], font_override_pt)
        f = p.runs[0].font if p.runs else p.font
        f.size = Pt(fnt_pt)
        f.bold = font_bold
        f.color.rgb = get_contrast_color(color_hex)  # 背景HEXから黒/白を自動選択
        p.alignment = font_align

    # タイトル（任意）
    if title:
        title_h = int(height * 0.12)
        tb = shapes.add_textbox(left, top - title_h, width, int(title_h * 0.9))
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = str(title)
        f = p.runs[0].font if p.runs else p.font
        f.size = Pt(_auto_font_pt(h_each[0], font_override_pt) + 4)
        f.bold = True
        p.alignment = PP_ALIGN.CENTER