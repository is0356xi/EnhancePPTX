# src/tools/renderers/plain_box.py
# -*- coding: utf-8 -*-
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Pt

from .. import utils
from ..schemas.plain_box import PlainBoxData

def render(slide, data: dict, geom: dict, context: dict):
    """
    シンプルなテキストボックス（枠線なし）を描画する。
    """
    # --- 1. スキーマによるバリデーションとデフォルト値の適用 ---
    try:
        model = PlainBoxData.model_validate(data)
        style = model.style
    except Exception as e:
        context["logger"].error(f"plain_box: データ形式が不正です. {e}")
        return

    # --- 2. 図形の追加と配置 ---
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        geom["left"],
        geom["top"],
        geom["width"],
        geom["height"],
    )

    # --- 3. デザインの設定 ---
    # 塗りつぶし
    shape.fill.solid()
    # utils が返す tuple を RGBColor オブジェクトに変換
    shape.fill.fore_color.rgb = RGBColor(*utils.hex_to_rgb(style.background_color)) # ★ 修正

    # 枠線を「なし」に設定
    shape.line.fill.background()

    # --- 4. テキストフレームの設定 ---
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = model.text
    p.font.size = Pt(style.font_size)
    # utils が返す tuple を RGBColor オブジェクトに変換
    p.font.color.rgb = RGBColor(*utils.hex_to_rgb(style.font_color)) # ★ 修正

    # 水平・垂直方向の配置
    align_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    p.alignment = align_map.get(style.align, PP_ALIGN.LEFT)

    valign_map = {
        "top": MSO_VERTICAL_ANCHOR.TOP,
        "middle": MSO_VERTICAL_ANCHOR.MIDDLE,
        "bottom": MSO_VERTICAL_ANCHOR.BOTTOM,
    }
    text_frame.vertical_anchor = valign_map.get(style.vertical_align, MSO_VERTICAL_ANCHOR.MIDDLE)