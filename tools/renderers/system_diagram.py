from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml import parse_xml

from ..utils import hex_to_rgb
from ..schemas.system_diagram import SystemDiagramSchema

# ノードタイプと図形のマッピング
NODE_SHAPES = {
    "user": MSO_SHAPE.ACTION_BUTTON_HOME,
    "system": MSO_SHAPE.ROUNDED_RECTANGLE,
    "database": MSO_SHAPE.CAN,
    "default": MSO_SHAPE.RECTANGLE,
}


def _calculate_connector_endpoints(from_geom, to_geom):
    """2つの図形のジオメトリから最適な接続点の座標を計算する"""
    fx, fy, fw, fh = from_geom["x"], from_geom["y"], from_geom["w"], from_geom["h"]
    tx, ty, tw, th = to_geom["x"], to_geom["y"], to_geom["w"], to_geom["h"]

    # 接続元の中心点
    from_cx = fx + fw / 2
    from_cy = fy + fh / 2
    # 接続先の中心点
    to_cx = tx + tw / 2
    to_cy = ty + th / 2

    # 水平方向か垂直方向か
    if abs(to_cx - from_cx) > abs(to_cy - from_cy):  # 水平方向の接続
        if from_cx < to_cx:  # 右へ
            start_point = (fx + fw, from_cy)
            end_point = (tx, to_cy)
        else:  # 左へ
            start_point = (fx, from_cy)
            end_point = (tx + tw, to_cy)
    else:  # 垂直方向の接続
        if from_cy < to_cy:  # 下へ
            start_point = (from_cx, fy + fh)
            end_point = (to_cx, ty)
        else:  # 上へ
            start_point = (from_cx, fy)
            end_point = (to_cx, ty + th)

    return start_point, end_point


def render(slide, data: SystemDiagramSchema, geom: dict, context: dict):
    """システム構成図を描画する"""
    left, top, width, height = geom["left"], geom["top"], geom["width"], geom["height"]

    grid = data.grid
    cell_width = width / grid.cols
    cell_height = height / grid.rows

    node_infos = {}

    # --- ノードのジオメトリを先に計算 ---
    for node_data in data.nodes:
        node_id = node_data.id
        pos = node_data.pos
        node_w = int(cell_width * 0.7)
        node_h = int(cell_height * 0.7)
        node_x = left + int(pos.col * cell_width + (cell_width - node_w) / 2)
        node_y = top + int(pos.row * cell_height + (cell_height - node_h) / 2)

        node_infos[node_id] = {
            "data": node_data,
            "pos": pos,
            "geom": {"x": node_x, "y": node_y, "w": node_w, "h": node_h},
        }

    # --- 境界線を先に描画 ---
    for boundary in data.boundaries:
        min_x, min_y, max_x, max_y = width * 2, height * 2, 0, 0
        for node_id in boundary.nodes:
            if node_id in node_infos:
                g = node_infos[node_id]["geom"]
                min_x, min_y = min(min_x, g["x"]), min(min_y, g["y"])
                max_x, max_y = max(max_x, g["x"] + g["w"]), max(max_y, g["y"] + g["h"])
        if max_x > min_x:
            padding = Pt(20)
            b_x, b_y = min_x - padding, min_y - padding
            b_w, b_h = (max_x - min_x) + 2 * padding, (max_y - min_y) + 2 * padding
            b_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, b_x, b_y, b_w, b_h
            )
            b_shape.fill.background()
            line = b_shape.line
            if boundary.color:
                line.color.rgb = RGBColor(*hex_to_rgb(boundary.color))
            if boundary.style == "dashed":
                line.dash_style = MSO_LINE_DASH_STYLE.DASH
            line.width = Pt(1.5)
            label_box = slide.shapes.add_textbox(
                b_x + Pt(5), b_y - Pt(15), Pt(200), Pt(20)
            )
            label_box.text_frame.text = boundary.label or ""
            label_box.text_frame.paragraphs[0].font.size = Pt(11)
            label_box.text_frame.paragraphs[0].font.italic = True

    # --- ノードを描画 ---
    for node_id, info in node_infos.items():
        node_type = info["data"].type or "default"
        shape_type = NODE_SHAPES.get(node_type, NODE_SHAPES["default"])
        g = info["geom"]
        shape = slide.shapes.add_shape(shape_type, g["x"], g["y"], g["w"], g["h"])
        info["shape"] = shape
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(230, 240, 255)
        shape.line.color.rgb = RGBColor(100, 120, 150)
        shape.line.width = Pt(1)
        tf = shape.text_frame
        tf.text = info["data"].label or ""
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True

    # --- コネクタを描画 ---
    for conn in data.connectors:
        from_node = node_infos.get(conn.from_)
        to_node = node_infos.get(conn.to)
        if not from_node or not to_node:
            continue

        # 修正: 始点と終点の座標を計算
        start_point, end_point = _calculate_connector_endpoints(
            from_node["geom"], to_node["geom"]
        )

        # 修正: 計算した座標でコネクタを追加
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            int(start_point[0]),
            int(start_point[1]),
            int(end_point[0]),
            int(end_point[1]),
        )

        # 矢印と線のスタイル設定 (この部分はXML操作で正しく動作)
        line = connector.line
        line_elm = line._get_or_add_ln()
        arrow_xml_ns = 'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        if conn.arrow_head in ["end", "both"]:
            headEnd = parse_xml(f'<a:headEnd type="arrow" {arrow_xml_ns}/>')
            line_elm.append(headEnd)
        if conn.arrow_head in ["start", "both"]:
            tailEnd = parse_xml(f'<a:tailEnd type="arrow" {arrow_xml_ns}/>')
            line_elm.append(tailEnd)
        if conn.style == "dashed":
            line.dash_style = MSO_LINE_DASH_STYLE.DASH

        if conn.label:
            mid_x = (start_point[0] + end_point[0]) / 2
            mid_y = (start_point[1] + end_point[1]) / 2
            label_box = slide.shapes.add_textbox(
                mid_x - Pt(40), mid_y - Pt(20), Pt(80), Pt(40)
            )
            tf = label_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = conn.label
            p.font.size = Pt(9)
            p.alignment = PP_ALIGN.CENTER