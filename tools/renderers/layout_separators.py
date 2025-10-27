from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from ..utils import hex_to_rgb
from ..schemas.layout_separators import LayoutSeparatorsSchema


def render(slide, data: LayoutSeparatorsSchema, geom: dict, context: dict):
    sections = data.sections
    if not sections:
        return

    num_sections = len(sections)
    section_width = geom["width"] / num_sections

    # --- ラベル描画 ---
    for i, label_text in enumerate(sections):
        label_x = geom["left"] + (i * section_width)

        tb = slide.shapes.add_textbox(
            label_x, geom["top"] - Pt(20), section_width, Pt(20)
        )
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label_text

        style = data.style
        r.font.size = Pt(style.font_size)
        r.font.bold = style.bold
        default_font_color = context.get("theme", {}).get("font_color", "#000000")
        font_color_hex = style.font_color or default_font_color
        r.font.color.rgb = RGBColor(*hex_to_rgb(font_color_hex))

    # --- 区切り線描画（Rectangleを細線に見せる方式） ---
    style = data.style
    line_color_hex = style.color

    for i in range(1, num_sections):
        line_x = geom["left"] + (i * section_width)

        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            line_x,
            geom["top"],
            Pt(1.0),  # 幅を細く
            geom["height"],
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(*hex_to_rgb(line_color_hex))
        line.line.fill.background()  # 枠線は消す