# -*- coding: utf-8 -*-
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Pt

from ..utils import hex_to_rgb
from ..schemas.weighted_process_timeline import WeightedProcessTimelineSchema


def render(slide, data: WeightedProcessTimelineSchema, geom: dict, context: dict):
    """工数付きプロセスタイムラインを描画する"""
    schema = data

    # タイトルの描画
    title_height = 0
    if schema.title:
        title_height = 30  # pt
        title_box = slide.shapes.add_textbox(
            geom["left"], geom["top"], geom["width"], Pt(title_height)
        )
        title_frame = title_box.text_frame
        title_frame.clear()
        title_frame.margin_left = 0
        title_frame.margin_right = 0
        title_frame.margin_top = 0
        title_frame.margin_bottom = 0

        p = title_frame.paragraphs[0]
        p.text = schema.title
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.bold = True

    # レーンの描画開始位置を計算
    content_top = geom["top"] + Pt(title_height)
    content_height = geom["height"] - Pt(title_height)

    if not schema.lanes:
        return

    # 全レーンの最大総重みを計算（正規化用）
    max_total_weight = 0
    for lane in schema.lanes:
        if lane.processes:
            lane_total_weight = sum(process.weight for process in lane.processes)
            max_total_weight = max(max_total_weight, lane_total_weight)

    if max_total_weight <= 0:
        return

    lane_count = len(schema.lanes)
    total_margin = Pt(schema.style.lane_margin_pt * (lane_count - 1))
    lane_height = (content_height - total_margin) / lane_count

    # 各レーンを描画
    current_y = content_top
    for lane_idx, lane in enumerate(schema.lanes):
        _render_lane(
            slide,
            lane,
            geom["left"],
            current_y,
            geom["width"],
            lane_height,
            schema.style,
            lane_idx,
            max_total_weight,
        )
        current_y += lane_height + Pt(schema.style.lane_margin_pt)


def _render_lane(
    slide,
    lane,
    lane_left,
    lane_top,
    lane_width,
    lane_height,
    style,
    lane_idx,
    max_total_weight,
):
    """単一レーンを描画する"""
    if not lane.processes:
        return

    # レーンの総重量を計算
    total_weight = sum(process.weight for process in lane.processes)
    if total_weight <= 0:
        return

    # レーン名ラベルの幅
    label_width = lane_width * 0.15  # レーン名用に15%確保
    process_area_width = lane_width - label_width

    # レーン名を描画
    if lane.name:
        label_box = slide.shapes.add_textbox(
            lane_left, lane_top, label_width, lane_height
        )
        label_frame = label_box.text_frame
        label_frame.clear()
        label_frame.margin_left = Pt(5)
        label_frame.margin_right = Pt(5)
        label_frame.margin_top = 0
        label_frame.margin_bottom = 0
        label_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = label_frame.paragraphs[0]
        p.text = lane.name
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(style.label_font_size)
        p.font.bold = True

    # プロセスを描画
    current_x = lane_left + label_width
    for process_idx, process in enumerate(lane.processes):
        # 最大総重みを基準に幅を計算（各レーンで統一されたスケール）
        process_width = (process.weight / max_total_weight) * process_area_width

        # マージンを考慮した実際の幅
        actual_width = process_width - Pt(style.process_margin_pt)
        if actual_width <= 0:
            current_x += process_width
            continue

        _render_process(
            slide,
            process,
            current_x,
            lane_top,
            actual_width,
            lane_height,
            style,
            lane_idx,
            process_idx,
            lane,
        )

        current_x += process_width


def _render_process(
    slide, process, x, y, width, height, style, lane_idx, process_idx, lane
):
    """単一プロセスを描画する"""
    # プロセスボックスを作成
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, width, height)

    # 色を設定（レーンの色を使用）
    process_color = _get_lane_color(lane, lane_idx)
    rgb = hex_to_rgb(process_color)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*rgb)

    # ボーダーを設定（枠線なし）
    shape.line.fill.background()

    # テキストを設定
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.margin_left = Pt(2)
    text_frame.margin_right = Pt(2)
    text_frame.margin_top = Pt(1)
    text_frame.margin_bottom = Pt(1)
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.word_wrap = True

    # プロセス名
    if style.show_labels and process.name:
        p = text_frame.paragraphs[0]
        p.text = process.name
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(style.label_font_size)
        p.font.bold = True

        # テキスト色を黒に固定
        p.font.color.rgb = RGBColor(0, 0, 0)

    # 重み表示
    if style.show_weights:
        if style.show_labels and process.name:
            # 新しい段落を追加
            p2 = text_frame.add_paragraph()
        else:
            p2 = text_frame.paragraphs[0]

        p2.text = f"({process.weight})"
        p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Pt(style.weight_font_size)

        # テキスト色を黒に固定
        p2.font.color.rgb = RGBColor(0, 0, 0)


def _get_lane_color(lane, lane_idx):
    """レーンの色を取得する"""
    if lane.color:
        return lane.color

    # 黒文字が読みやすい薄いカラーパレット
    light_colors = [
        "#E3F2FD",  # 非常に薄い青
        "#E8F5E8",  # 非常に薄い緑
        "#FFF9E1",  # 非常に薄い黄色
        "#FFEBEE",  # 非常に薄いピンク
        "#F3E5F5",  # 非常に薄い紫
        "#FFF3E0",  # 非常に薄いオレンジ
    ]

    # レーンインデックスに基づいて色を決定
    color_idx = lane_idx % len(light_colors)
    return light_colors[color_idx]
