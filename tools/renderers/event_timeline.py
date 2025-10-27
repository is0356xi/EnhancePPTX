from __future__ import annotations

from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml import parse_xml

from ..utils import (
    hex_to_rgb,
    normalize_color,
    get_contrast_color,
)
from ..schemas.event_timeline import EventTimelineSchema


def render(slide, data: EventTimelineSchema, geom: dict, context: dict):
    """横並びイベントタイムラインを描画するレンダラ。

    - 上部に横線と番号付きサークル
    - サークルの下に色付きバナー（小見出し）
    - さらに下に白いカード（詳細）を並べる
    """
    shapes = (context or {}).get("shapes_target") or slide.shapes

    left = geom["left"]
    top = geom["top"]
    width = geom["width"]
    height = geom["height"]

    steps = list(data.steps or [])
    n = max(1, len(steps))

    # 基本寸法
    line_y = top + int(height * 0.12)
    circle_d = int(min(width / (n * 6), height * 0.10))
    circle_y = line_y - int(circle_d / 2)

    # 横線を細い右向きコネクタ矢印に置換
    try:
        line_h = int(Pt(2))
        line_left = left + int(Pt(6))
        line_right = left + width - int(Pt(6))
        start_x = line_left
        start_y = line_y
        end_x = line_right
        end_y = line_y

        # 直線コネクタを追加して矢印ヘッドを付ける
        connector = shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, int(start_x), int(start_y), int(end_x), int(end_y)
        )
        line = connector.line
        line.width = Pt(1.5)
        line.color.rgb = RGBColor(180, 180, 180)

        # XML で arrow head を追加
        try:
            line_elm = line._get_or_add_ln()
            arrow_xml_ns = (
                'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
            )
            # attach arrow at the start (tail) so the arrow points left
            try:
                tailEnd = parse_xml(f'<a:tailEnd type="arrow" {arrow_xml_ns}/>')
                line_elm.append(tailEnd)
            except Exception:
                # ignore XML failure
                pass
        except Exception:
            # フォールバック
            line_left = left
            line_w = width
    except Exception:
        # Fallback if connector creation failed entirely
        line_left = left
        line_w = width

    # 横方向の中心位置配分
    gap = int((width - 2 * int(Pt(6))) / n)

    for i, step in enumerate(steps):
        cx = left + int(Pt(6)) + int(gap * (i + 0.5)) - int(circle_d / 2)

        # サークル（番号表示）
        try:
            circ = shapes.add_shape(MSO_SHAPE.OVAL, cx, circle_y, circle_d, circle_d)
            circ.fill.solid()
            r, g, b = normalize_color(step.color or data.card_bg)
            circ.fill.fore_color.rgb = RGBColor(r, g, b)
            circ.line.fill.background()
            if data.show_numbers:
                tf = circ.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = str(i + 1)
                p.level = 0
                r2, g2, b2 = get_contrast_color(step.color or data.card_bg)
                p.runs[0].font.size = Pt(12)
                p.runs[0].font.bold = True
                p.runs[0].font.color.rgb = r2
        except Exception:
            pass

        # 一貫した幅を使う（バナーとカードを同じ幅にして揃える）
        shared_w = int(gap * 0.95)

        # バナー（色帯・年など短いテキスト）
        try:
            banner_w = shared_w
            banner_h = int(height * 0.18)
            banner_top = circle_y + circle_d + int(Pt(6))
            bx = cx - int((banner_w - circle_d) / 2)
            # Use plain rectangle (no rounded corners) so banner looks square
            banner = shapes.add_shape(
                MSO_SHAPE.RECTANGLE, bx, banner_top, banner_w, banner_h
            )
            try:
                r, g, b = normalize_color(step.color or "#0D6EFD")
                banner.fill.solid()
                banner.fill.fore_color.rgb = RGBColor(r, g, b)
                banner.line.fill.background()
            except Exception:
                pass
            # 右下に小さな年テキストを淡色で入れる
            if step.badge:
                try:
                    tf = banner.text_frame
                    tf.clear()
                    tf.word_wrap = True
                    # small margin
                    try:
                        tf.margin_left = tf.margin_right = tf.margin_top = (
                            tf.margin_bottom
                        ) = 0
                    except Exception:
                        pass
                    p = tf.paragraphs[0]
                    p.text = str(step.badge)
                    p.level = 0
                    p.runs[0].font.size = Pt(14)
                    p.runs[0].font.bold = True
                    p.runs[0].font.color.rgb = RGBColor(255, 255, 255)
                except Exception:
                    pass
        except Exception:
            pass

        # カード（白）
        try:
            card_w = shared_w
            card_h = int(height * 0.36)
            card_top = banner_top + banner_h + int(Pt(8))
            card_left = bx
            # Use plain rectangle (no rounded corners) so card looks square
            card = shapes.add_shape(
                MSO_SHAPE.RECTANGLE, card_left, card_top, card_w, card_h
            )
            try:
                r_bg, g_bg, b_bg = normalize_color(data.card_bg)
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(r_bg, g_bg, b_bg)
                # default no border; border may be set below
            except Exception:
                pass
            # 枠線風に細い線を同色で描画（薄いグレー）
            try:
                r, g, b = normalize_color(data.card_border)
                card.line.color.rgb = RGBColor(r, g, b)
            except Exception:
                pass

            # カード内にステップの説明を入れる（図形内テキスト）
            # - description を主テキストとして使う
            description = step.description or ""
            try:
                tfc = card.text_frame
                tfc.clear()
                tfc.word_wrap = True
                try:
                    tfc.margin_left = tfc.margin_right = tfc.margin_top = (
                        tfc.margin_bottom
                    ) = 0
                except Exception:
                    pass

                p = tfc.paragraphs[0]
                if description:
                    p.text = description
                    p.level = 0
                    p.runs[0].font.size = Pt(12)
                    p.runs[0].font.bold = False
                    p.runs[0].font.color.rgb = RGBColor(51, 51, 51)
            except Exception:
                pass
        except Exception:
            pass
